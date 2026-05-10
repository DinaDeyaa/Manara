from __future__ import annotations

import hashlib
import os
import io
import re
import json
import base64
import zipfile
import random
import threading
import time
from pathlib import Path
from collections import Counter
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Callable, Optional

import fitz
import pandas as pd
import chromadb
from tqdm import tqdm
from openai import OpenAI, RateLimitError
from PIL import Image
from pptx import Presentation
from chromadb.utils import embedding_functions
from pptx.enum.shapes import MSO_SHAPE_TYPE


# =========================================================
# CONFIG
# =========================================================

PROJECT_DIR = Path("/Users/dinaal-memah/Desktop/graduation project 2")
COURSES_DIR = PROJECT_DIR / "courses"
FAILED_VISION_PAGES_FILE = PROJECT_DIR / "failed_vision_pages.json"

# Tier-1 preprocessing models with stable public limits.
MODEL_NAME        = "gpt-5-nano"      # text: summaries, concept extraction
VISION_MODEL_NAME = "gpt-4.1-mini"    # vision: images, diagrams, slides

# OPT: reduced token limits — concept JSON rarely needs more than 500 tokens,
#      summaries rarely exceed 600. Avoids over-allocating output buffer.
MAX_OUTPUT_TOKENS        = 900   # (summaries)
MAX_CONCEPT_OUTPUT_TOKENS = 900  # (concepts)
MAX_VISION_OUTPUT_TOKENS = 500  # constrained for low TPM account tiers

# GPT-5 chat completions count hidden reasoning tokens inside
# max_completion_tokens. Keep reasoning tiny and reserve enough room for the
# visible summary / JSON payload so semantic extraction does not return "".
TEXT_LLM_REASONING_EFFORT = "minimal"
MIN_TEXT_COMPLETION_BUDGET = 1800

RUN_SUMMARIZATION    = True
RUN_CONCEPT_EXTRACTION = True
SAVE_TO_CHROMA       = True

FORCE_REGENERATE = False
TEST_COURSE_NAME = False

# Temporary semantic-stage diagnostics. Keep enabled while investigating empty
# summaries/concepts, then set to False after the pipeline is healthy.
DEBUG_SEMANTIC_EXTRACTION = True
DEBUG_LLM_PREVIEW_CHARS   = 1200

# Vision controls
ENABLE_VISION                    = True
VISION_QUOTA_EXHAUSTED          = False
RUN_VISION_EXTRACTION           = True
VISION_ONLY_WHEN_NEEDED         = True
MIN_TEXT_LENGTH_FOR_VISION_SKIP = 150
IMAGE_TEXT_LENGTH_FOR_VISION    = 300
MAX_VISION_PAGES_PER_DOCUMENT   = 3
MAX_TOTAL_VISION_CALLS          = 40
MAX_VISION_IMAGE_WIDTH          = 900
VISION_JPEG_QUALITY             = 55
MAX_VISION_CHARS_PER_PAGE       = 4000  # was 6000

# OPT: parallelism — all IO-bound so threading gives near-linear speedup
MAX_PAGE_WORKERS   = 6   # threads per document for page/slide vision calls
MAX_FILE_WORKERS   = 4   # threads for summarization + concepts across files
MAX_COURSE_WORKERS = 2   # concurrent courses

# Rate-limit controls. File/page workers may still run in parallel, but only a
# bounded number of image requests can enter the OpenAI vision endpoint globally.
OPENAI_MAX_RETRIES = 5
OPENAI_RETRY_BACKOFFS = [1, 2, 4, 8, 16]
VISION_SUCCESS_SLEEP_SECONDS = 2
TEXT_SUCCESS_SLEEP_SECONDS = 1
MAX_GLOBAL_VISION_REQUESTS = 3
VISION_SEMAPHORE = threading.Semaphore(3)

# OPT: reduced PDF render zoom for low TPM account tiers.
PDF_RENDER_ZOOM = 1.0

SUPPORTED_DIRECT_FILE_EXTENSIONS = {
    ".pdf",
    ".pptx",
    ".cpp",
    ".h",
    ".c",
    ".hpp",
    ".py",
    ".java",
    ".js",
    ".ts",
    ".txt",
    ".md",
    ".ipynb",
}

SUPPORTED_ZIP_MEMBER_EXTENSIONS = SUPPORTED_DIRECT_FILE_EXTENSIONS

IGNORED_DIR_NAMES = {
    "outputs",
    "chroma_db",
    "__pycache__",
    ".git",
    ".idea",
    ".vscode",
    "__MACOSX",
}

CHUNKS_COLLECTION_SUFFIX = "_chunks"

# OPT: lazy client init — creating OpenAI() at module level crashes if
#      OPENAI_API_KEY is not set (e.g. during Uvicorn reload). Now created
#      on first actual LLM call only.
_openai_client: Optional[OpenAI] = None

def get_client() -> OpenAI:
    global _openai_client
    if _openai_client is None:
        key = os.getenv("OPENAI_API_KEY")
        if not key:
            raise RuntimeError("OPENAI_API_KEY is not set in your environment.")
        _openai_client = OpenAI(api_key=key)
    return _openai_client


_vision_stats_lock = threading.Lock()
_vision_budget_exhausted = threading.Event()
_openai_quota_exhausted = threading.Event()
VISION_STATS = {
    "vision_total_calls": 0,
    "vision_success": 0,
    "vision_retries": 0,
    "vision_failures": 0,
    "vision_pages_skipped": 0,
    "vision_latency_total_seconds": 0.0,
    "vision_estimated_tokens": 0,
    "vision_pages_skipped_quota": 0,
    "ocr_recoveries": 0,
    "total_files": 0,
    "text_only_pages_processed": 0,
    "total_retries": 0,
    "total_wait_time_seconds": 0.0,
    "current_tpm_wait_events": 0,
}

_failed_vision_lock = threading.Lock()
_failed_vision_pages: list[dict] = []


def _increment_vision_stat(key: str, amount: int | float = 1) -> None:
    with _vision_stats_lock:
        VISION_STATS[key] = VISION_STATS.get(key, 0) + amount


def _tracked_sleep(seconds: float) -> None:
    _increment_vision_stat("total_wait_time_seconds", seconds)
    time.sleep(seconds)


def debug_semantic(message: str) -> None:
    if DEBUG_SEMANTIC_EXTRACTION:
        print(f"[SEMANTIC DEBUG] {message}")


def debug_preview(value, limit: int = DEBUG_LLM_PREVIEW_CHARS) -> str:
    text = str(value)
    text = text.replace("\n", "\\n")
    if len(text) > limit:
        return text[:limit] + "...<truncated>"
    return text


def count_concept_items(concepts: dict | None) -> tuple[int, int]:
    if not isinstance(concepts, dict):
        return 0, 0
    topic_count = 0
    subtopic_count = 0
    for content in concepts.values():
        if not isinstance(content, dict):
            continue
        topics = content.get("topics", [])
        if not isinstance(topics, list):
            continue
        topic_count += len(topics)
        for topic in topics:
            if isinstance(topic, dict):
                subtopics = topic.get("subtopics", [])
                if isinstance(subtopics, list):
                    subtopic_count += len([s for s in subtopics if str(s).strip()])
    return topic_count, subtopic_count


def print_vision_counters(prefix: str = "  Vision counters") -> None:
    with _vision_stats_lock:
        snapshot = dict(VISION_STATS)
    total_calls = snapshot.get("vision_total_calls", 0)
    avg_latency = (
        snapshot.get("vision_latency_total_seconds", 0.0) / total_calls
        if total_calls
        else 0.0
    )
    print(
        f"{prefix}: total_calls={total_calls}, "
        f"avg_latency={avg_latency:.2f}s, "
        f"success={snapshot['vision_success']}, "
        f"pages_skipped={snapshot['vision_pages_skipped']}, "
        f"quota_skipped={snapshot['vision_pages_skipped_quota']}, "
        f"text_only_pages={snapshot['text_only_pages_processed']}, "
        f"ocr_recoveries={snapshot['ocr_recoveries']}, "
        f"estimated_tokens={snapshot['vision_estimated_tokens']}, "
        f"pages_retried={snapshot['vision_retries']}, "
        f"total_retries={snapshot['total_retries']}, "
        f"total_wait_time={snapshot['total_wait_time_seconds']:.2f}s, "
        f"failures={snapshot['vision_failures']}, "
        f"tpm_wait_events={snapshot['current_tpm_wait_events']}"
    )


def print_preprocessing_statistics() -> None:
    with _vision_stats_lock:
        snapshot = dict(VISION_STATS)
    print(
        "[PREPROCESSING STATS] "
        f"total_files={snapshot['total_files']}, "
        f"total_vision_calls_completed={snapshot['vision_success']}, "
        f"total_vision_pages_skipped_due_to_quota={snapshot['vision_pages_skipped_quota']}, "
        f"total_text_only_pages_processed={snapshot['text_only_pages_processed']}, "
        f"total_ocr_recoveries={snapshot['ocr_recoveries']}, "
        f"estimated_token_usage={snapshot['vision_estimated_tokens']}"
    )


def print_startup_rate_limit_config() -> None:
    print(
        "[PREPROCESSING CONFIG] "
        f"model={MODEL_NAME}, "
        f"vision_model={VISION_MODEL_NAME}, "
        f"max_file_workers={MAX_FILE_WORKERS}, "
        f"max_page_workers={MAX_PAGE_WORKERS}, "
        f"max_course_workers={MAX_COURSE_WORKERS}, "
        f"vision_semaphore_size={MAX_GLOBAL_VISION_REQUESTS}, "
        f"text_reasoning_effort={TEXT_LLM_REASONING_EFFORT}, "
        f"min_text_completion_budget={MIN_TEXT_COMPLETION_BUDGET}, "
        f"pdf_render_zoom={PDF_RENDER_ZOOM}, "
        f"max_vision_output_tokens={MAX_VISION_OUTPUT_TOKENS}, "
        f"enable_vision={ENABLE_VISION}, "
        f"min_text_skip={MIN_TEXT_LENGTH_FOR_VISION_SKIP}, "
        f"image_text_threshold={IMAGE_TEXT_LENGTH_FOR_VISION}, "
        f"max_vision_pages_per_document={MAX_VISION_PAGES_PER_DOCUMENT}, "
        f"max_total_vision_calls={MAX_TOTAL_VISION_CALLS}, "
        f"max_image_width={MAX_VISION_IMAGE_WIDTH}, "
        f"jpeg_quality={VISION_JPEG_QUALITY}, "
        f"vision_success_sleep={VISION_SUCCESS_SLEEP_SECONDS}s, "
        f"text_success_sleep={TEXT_SUCCESS_SLEEP_SECONDS}s"
    )


def _is_rate_limit_error(exc: Exception) -> bool:
    status_code = getattr(exc, "status_code", None)
    code = getattr(exc, "code", None)
    message = str(exc).lower()
    return (
        isinstance(exc, RateLimitError)
        or status_code == 429
        or code == "rate_limit_exceeded"
        or "rate_limit" in message
        or "rate limit" in message
    )


def _is_insufficient_quota_error(exc: Exception) -> bool:
    code = getattr(exc, "code", None)
    body = getattr(exc, "body", None)
    if isinstance(body, dict):
        error = body.get("error", {})
        if isinstance(error, dict):
            code = code or error.get("code")
    message = str(exc).lower()
    return code == "insufficient_quota" or "insufficient_quota" in message


def _is_transient_network_error(exc: Exception) -> bool:
    status_code = getattr(exc, "status_code", None)
    message = str(exc).lower()
    return (
        status_code in {502, 503, 504}
        or "timeout" in message
        or "timed out" in message
        or "connection" in message
        or "temporarily unavailable" in message
        or "server error" in message
    )


def _disable_vision_for_quota() -> None:
    global ENABLE_VISION, VISION_QUOTA_EXHAUSTED
    if not _openai_quota_exhausted.is_set():
        _openai_quota_exhausted.set()
        VISION_QUOTA_EXHAUSTED = True
        ENABLE_VISION = False
        print("[OPENAI QUOTA EXHAUSTED] Vision API quota unavailable. Continuing with text-only extraction.")


def _retry_after_seconds(exc: Exception) -> float | None:
    message = str(exc)
    match = re.search(r"try again in\s+([0-9]+(?:\.[0-9]+)?)s", message, flags=re.IGNORECASE)
    if match:
        return float(match.group(1))
    return None


def _handle_rate_limit_wait(exc: Exception, attempt: int, *, is_vision: bool) -> None:
    backoff = OPENAI_RETRY_BACKOFFS[min(attempt, len(OPENAI_RETRY_BACKOFFS) - 1)]
    retry_after = _retry_after_seconds(exc)
    wait_seconds = backoff + random.uniform(0, 0.5)
    if retry_after is not None:
        wait_seconds = max(wait_seconds, retry_after + 0.5)

    _increment_vision_stat("total_retries")

    if is_vision:
        _increment_vision_stat("vision_retries")
        _increment_vision_stat("current_tpm_wait_events")

    print(
        f"  OpenAI rate limit hit; retry {attempt + 1}/{OPENAI_MAX_RETRIES} "
        f"in {wait_seconds:.2f}s"
    )
    try:
        _tracked_sleep(wait_seconds)
    finally:
        if is_vision:
            _increment_vision_stat("current_tpm_wait_events", -1)


def _call_openai_with_retries(call_fn: Callable, *, is_vision: bool = False):
    """
    Retry OpenAI calls that hit TPM/rate-limit pressure. Non-rate-limit errors
    still fail fast so malformed requests and auth issues do not get hidden.
    """
    last_error = None

    for attempt in range(OPENAI_MAX_RETRIES + 1):
        if is_vision:
            with VISION_SEMAPHORE:
                with _vision_stats_lock:
                    vision_calls_used = VISION_STATS.get("vision_total_calls", 0)
                if vision_calls_used >= MAX_TOTAL_VISION_CALLS:
                    _vision_budget_exhausted.set()
                    _increment_vision_stat("vision_pages_skipped")
                    print("  Vision daily safety cap reached; disabling further vision calls for this run.")
                    return None

                start = time.time()
                try:
                    result = call_fn()
                    elapsed = time.time() - start
                    _increment_vision_stat("vision_total_calls")
                    _increment_vision_stat("vision_latency_total_seconds", elapsed)
                    _increment_vision_stat("vision_success")
                    _tracked_sleep(VISION_SUCCESS_SLEEP_SECONDS)
                    return result
                except Exception as exc:
                    elapsed = time.time() - start
                    _increment_vision_stat("vision_total_calls")
                    _increment_vision_stat("vision_latency_total_seconds", elapsed)
                    last_error = exc
                    if _is_insufficient_quota_error(exc):
                        _disable_vision_for_quota()
                        _increment_vision_stat("vision_pages_skipped_quota")
                        return None
                    if (
                        not (_is_rate_limit_error(exc) or _is_transient_network_error(exc))
                        or attempt >= OPENAI_MAX_RETRIES
                    ):
                        raise
                    _handle_rate_limit_wait(exc, attempt, is_vision=True)
                    continue

        try:
            result = call_fn()
            _tracked_sleep(TEXT_SUCCESS_SLEEP_SECONDS)
            return result
        except Exception as exc:
            last_error = exc
            if _is_insufficient_quota_error(exc):
                raise
            if (
                not (_is_rate_limit_error(exc) or _is_transient_network_error(exc))
                or attempt >= OPENAI_MAX_RETRIES
            ):
                raise
            _handle_rate_limit_wait(exc, attempt, is_vision=False)
            continue

    raise last_error


def _load_failed_vision_pages() -> list[dict]:
    if not FAILED_VISION_PAGES_FILE.exists():
        return []
    try:
        with open(FAILED_VISION_PAGES_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, list) else []
    except Exception:
        return []


def _write_failed_vision_pages(entries: list[dict]) -> None:
    FAILED_VISION_PAGES_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(FAILED_VISION_PAGES_FILE, "w", encoding="utf-8") as f:
        json.dump(entries, f, indent=2, ensure_ascii=False)


def record_failed_vision_page(entry: dict) -> None:
    """Persist enough information to retry the exact page/slide later."""
    clean_entry = {
        "course_dir": entry.get("course_dir", ""),
        "course_name": entry.get("course_name", ""),
        "file_name": entry.get("file_name", ""),
        "relative_path": entry.get("relative_path", ""),
        "file_type": entry.get("file_type", ""),
        "page": entry.get("page"),
        "error": entry.get("error", ""),
    }

    with _failed_vision_lock:
        _failed_vision_pages.append(clean_entry)
        existing = _load_failed_vision_pages()
        existing.append(clean_entry)
        _write_failed_vision_pages(existing)

    _increment_vision_stat("vision_failures")


# =========================================================
# COURSE METADATA MAP
# =========================================================

COURSE_METADATA_MAP = {

    # ── Year 1 ────────────────────────────────────────────────────────────────

    "introduction to computer science": {
        "course_code": "11102",
        "official_title": "Introduction to Computer Science",
        "credit_hours": 3,
        "prerequisites": [],
        "concurrent": [],
        "recommended_year": 1,
        "recommended_semester": 1,
        "description": (
            "Introduces core computer science concepts, data representation, number systems, "
            "problem solving, flowcharts, and basic programming."
        ),
        "aliases": [],
    },
    "calculus 1": {
        "course_code": "20132",
        "official_title": "Calculus (1)",
        "credit_hours": 3,
        "prerequisites": [],
        "concurrent": [],
        "recommended_year": 1,
        "recommended_semester": 1,
        "description": (
            "Functions, limits and continuity, derivatives, differentiation, inverse functions, "
            "trigonometric functions, logarithmic and exponential functions, hyperbolic functions, and integrals."
        ),
        "aliases": ["calculus (1)", "calculus i"],
    },
    "discrete mathematics 1": {
        "course_code": "20134",
        "official_title": "Discrete Mathematics (1)",
        "credit_hours": 3,
        "prerequisites": [],
        "concurrent": [],
        "recommended_year": 1,
        "recommended_semester": 1,
        "description": (
            "Introduces discrete mathematical structures including logic, sets, relations, functions, "
            "mathematical induction, counting, graphs, and trees, with applications to computer science."
        ),
        "aliases": ["discrete mathematics (1)", "discrete math 1", "discrete math (1)"],
    },
    "structured programming": {
        "course_code": "11103",
        "official_title": "Structured Programming",
        "credit_hours": 3,
        "prerequisites": ["Introduction to Computer Science"],
        "concurrent": ["Structured Programming Lab"],
        "recommended_year": 1,
        "recommended_semester": 2,
        "description": (
            "Introduces structured programming concepts, C++ syntax and semantics, control "
            "structures, recursion, functions, arrays, pointers, and basic file I/O."
        ),
        "aliases": [],
    },
    "structured programming lab": {
        "course_code": "11151",
        "official_title": "Structured Programming Lab",
        "credit_hours": 1,
        "prerequisites": ["Structured Programming"],
        "concurrent": ["Structured Programming"],
        "recommended_year": 1,
        "recommended_semester": 2,
        "description": (
            "Practical lab exercises complementing Structured Programming, covering C++ "
            "implementation of control structures, functions, arrays, and file I/O."
        ),
        "aliases": [],
    },
    "calculus 2": {
        "course_code": "20133",
        "official_title": "Calculus (2)",
        "credit_hours": 3,
        "prerequisites": ["Calculus 1"],
        "concurrent": [],
        "recommended_year": 1,
        "recommended_semester": 2,
        "description": (
            "Methods of integration, applications of integration, plane analytic geometry including "
            "polar coordinates, sequences and series, including power series."
        ),
        "aliases": ["calculus (2)", "calculus ii"],
    },
    "introduction to data science": {
        "course_code": "14140",
        "official_title": "Introduction to Data Science",
        "credit_hours": 3,
        "prerequisites": ["Introduction to Computer Science"],
        "concurrent": [],
        "recommended_year": 1,
        "recommended_semester": 2,
        "description": (
            "Introduces students to the field of data science and its basic principles, "
            "tools, data collection and integration, exploratory data analysis, predictive "
            "and descriptive modeling, evaluation, and communication."
        ),
        "aliases": [],
    },
    "statistical methods": {
        "course_code": "20233",
        "official_title": "Statistical Methods",
        "credit_hours": 3,
        "prerequisites": [],
        "concurrent": [],
        "recommended_year": 1,
        "recommended_semester": 2,
        "description": (
            "Covers descriptive statistics, probability distributions, statistical inference, "
            "hypothesis testing, regression analysis, and applications in data analysis."
        ),
        "aliases": [],
    },
    "technical writing and communication skills": {
        "course_code": "20200",
        "official_title": "Technical Writing and Communication Skills",
        "credit_hours": 3,
        "prerequisites": [],
        "concurrent": [],
        "recommended_year": 1,
        "recommended_semester": 2,
        "description": (
            "Develops technical writing and professional communication skills including report "
            "writing, documentation, presentations, and scientific communication for engineers "
            "and computer scientists."
        ),
        "aliases": ["technical writing"],
    },

    # ── Year 2 ────────────────────────────────────────────────────────────────

    "object oriented programming": {
        "course_code": "11206",
        "official_title": "Object Oriented Programming",
        "credit_hours": 3,
        "prerequisites": ["Structured Programming"],
        "concurrent": ["Object Oriented Programming Lab"],
        "recommended_year": 2,
        "recommended_semester": 1,
        "description": (
            "Introduces OOP concepts including abstraction, encapsulation, classes, inheritance, "
            "overloading, polymorphism, and templates."
        ),
        "aliases": ["oop"],
    },
    "object oriented programming lab": {
        "course_code": "11253",
        "official_title": "Object Oriented Programming Lab",
        "credit_hours": 1,
        "prerequisites": ["Object Oriented Programming"],
        "concurrent": ["Object Oriented Programming"],
        "recommended_year": 2,
        "recommended_semester": 1,
        "description": (
            "Practical lab exercises complementing Object Oriented Programming, covering "
            "implementation of classes, inheritance, polymorphism, and design patterns in C++ or Java."
        ),
        "aliases": ["oop lab"],
    },
    "linear algebra": {
        "course_code": "20234",
        "official_title": "Linear Algebra",
        "credit_hours": 3,
        "prerequisites": ["Calculus 2"],
        "concurrent": [],
        "recommended_year": 2,
        "recommended_semester": 1,
        "description": (
            "System of linear equations, row-echelon form, Gaussian elimination, Gauss-Jordan method, "
            "matrices, determinants, Euclidean n-space, linear transformations, vector spaces, "
            "orthogonality, least squares, QR decomposition, eigenvalues and eigenvectors."
        ),
        "aliases": [],
    },
    "data structures and introduction to algorithms": {
        "course_code": "11212",
        "official_title": "Data Structures and Introduction to Algorithms",
        "credit_hours": 3,
        "prerequisites": ["Object Oriented Programming"],
        "concurrent": [],
        "recommended_year": 2,
        "recommended_semester": 2,
        "description": (
            "Introduces algorithm design and analysis basics, asymptotic complexity, searching, "
            "sorting, recursion, and core data structures."
        ),
        "aliases": ["data structures", "dsa"],
    },
    "data engineering": {
        "course_code": "14260",
        "official_title": "Data Engineering",
        "credit_hours": 3,
        "prerequisites": ["Introduction to Data Science"],
        "concurrent": ["Data Engineering Lab"],
        "recommended_year": 2,
        "recommended_semester": 1,
        "description": (
            "Examines the modern data ecosystem and core ETL tasks, data types, staging, "
            "profiling, cleansing, migration, and basic data visualization."
        ),
        "aliases": [],
    },
    "data engineering lab": {
        "course_code": "14261",
        "official_title": "Data Engineering Lab",
        "credit_hours": 1,
        "prerequisites": ["Data Engineering"],
        "concurrent": ["Data Engineering"],
        "recommended_year": 2,
        "recommended_semester": 1,
        "description": (
            "Practical exercises using common data engineering tools and ETL tasks on "
            "different types of data."
        ),
        "aliases": [],
    },
    "statistics and probability for data science": {
        "course_code": "14270",
        "official_title": "Statistics and Probability for Data Science",
        "credit_hours": 3,
        "prerequisites": ["Statistical Methods"],
        "concurrent": [],
        "recommended_year": 2,
        "recommended_semester": 2,
        "description": (
            "Covers probability theory, random variables, probability distributions, statistical "
            "inference, Bayesian methods, and their applications in data science and machine learning."
        ),
        "aliases": ["stats and probability", "probability for data science"],
    },
    "high performance computing for big data": {
        "course_code": "14362",
        "official_title": "High Performance Computing for Big Data",
        "credit_hours": 3,
        "prerequisites": ["Data Engineering"],
        "concurrent": [],
        "recommended_year": 2,
        "recommended_semester": 2,
        "description": (
            "Introduces big data concepts, platform organization, tools such as Hadoop and "
            "Spark, and the workflow of big-data components."
        ),
        "aliases": ["hpc", "big data"],
    },

    # ── Year 3 ────────────────────────────────────────────────────────────────

    "algorithms design and analysis": {
        "course_code": "11313",
        "official_title": "Algorithms Design and Analysis",
        "credit_hours": 3,
        "prerequisites": ["Data Structures and Introduction to Algorithms"],
        "concurrent": [],
        "recommended_year": 3,
        "recommended_semester": 1,
        "description": (
            "Covers formal techniques for designing and analyzing algorithms including greedy, "
            "divide-and-conquer, backtracking, heuristics, and graph algorithms."
        ),
        "aliases": ["algorithms", "ada"],
    },
    "database systems": {
        "course_code": "11323",
        "official_title": "Database Systems",
        "credit_hours": 3,
        "prerequisites": ["Data Structures and Introduction to Algorithms"],
        "concurrent": ["Database Systems Lab"],
        "recommended_year": 3,
        "recommended_semester": 1,
        "description": (
            "Covers basic database concepts, DBMS components, transaction management, "
            "data modeling, ER diagrams, relational algebra, queries, and normalization."
        ),
        "aliases": ["db", "dbms"],
    },
    "database systems lab": {
        "course_code": "11354",
        "official_title": "Database Systems Lab",
        "credit_hours": 1,
        "prerequisites": ["Database Systems"],
        "concurrent": ["Database Systems"],
        "recommended_year": 3,
        "recommended_semester": 1,
        "description": (
            "Practical training on designing and implementing a full database application using "
            "a relational DBMS."
        ),
        "aliases": ["db lab"],
    },
    "operating systems": {
        "course_code": "11335",
        "official_title": "Operating Systems",
        "credit_hours": 3,
        "prerequisites": ["Data Structures and Introduction to Algorithms"],
        "concurrent": ["Operating Systems Lab"],
        "recommended_year": 3,
        "recommended_semester": 1,
        "description": (
            "Introduction to operating systems, processes, threads, CPU scheduling, process synchronization, "
            "deadlocks, memory management, virtual memory, file systems, mass storage management, and UNIX case study."
        ),
        "aliases": ["os"],
    },
    "operating systems lab": {
        "course_code": "11355",
        "official_title": "Operating Systems Lab",
        "credit_hours": 1,
        "prerequisites": ["Operating Systems"],
        "concurrent": ["Operating Systems"],
        "recommended_year": 3,
        "recommended_semester": 1,
        "description": (
            "Practical UNIX/Linux skills including installation, Vi editor, file and process management, "
            "shell programming, system administration, and implementation of some operating system concepts."
        ),
        "aliases": ["os lab"],
    },
    "artificial intelligence": {
        "course_code": "14330",
        "official_title": "Artificial Intelligence",
        "credit_hours": 3,
        "prerequisites": ["Data Structures and Introduction to Algorithms"],
        "concurrent": [],
        "recommended_year": 3,
        "recommended_semester": 1,
        "description": (
            "Introduces AI concepts, knowledge representation, heuristic search, expert systems, "
            "natural language processing, machine learning, and AI applications."
        ),
        "aliases": ["ai"],
    },
    "data visualization": {
        "course_code": "14364",
        "official_title": "Data Visualization",
        "credit_hours": 3,
        "prerequisites": ["High Performance Computing for Big Data"],
        "concurrent": [],
        "recommended_year": 3,
        "recommended_semester": 2,
        "description": (
            "Covers designing and creating data visualizations, visual encoding, dashboard "
            "development, and identification of patterns and trends."
        ),
        "aliases": [],
    },
    "data mining": {
        "course_code": "14465",
        "official_title": "Data Mining",
        "credit_hours": 3,
        "prerequisites": ["High Performance Computing for Big Data"],
        "concurrent": [],
        "recommended_year": 3,
        "recommended_semester": 2,
        "description": (
            "Introduces data mining concepts and methods with focus on pattern discovery, "
            "clustering, classification, and anomaly detection."
        ),
        "aliases": [],
    },

    # ── Year 4 ────────────────────────────────────────────────────────────────

    "distributed systems": {
        "course_code": "11356",
        "official_title": "Distributed Systems",
        "credit_hours": 3,
        "prerequisites": ["Operating Systems"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Concepts of distributed systems, communication, client-server model, RPC, RMI, group communication, "
            "synchronization, election algorithms, atomic transactions, deadlocks, allocation, scheduling, "
            "fault tolerance, real-time systems, and distributed shared memory."
        ),
        "aliases": [],
    },
    "computer and society": {
        "course_code": "11449",
        "official_title": "Computer and Society",
        "credit_hours": 3,
        "prerequisites": [],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Explores the social, ethical, legal, and professional issues related to computing, "
            "including privacy, intellectual property, security, professional responsibility, "
            "and the societal impact of emerging technologies."
        ),
        "aliases": ["computers and society"],
    },
    "information systems security": {
        "course_code": "11464",
        "official_title": "Information Systems Security",
        "credit_hours": 3,
        "prerequisites": ["Operating Systems"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Covers fundamentals of information security including cryptography, authentication, "
            "access control, network security, vulnerability assessment, and security policies."
        ),
        "aliases": ["information security", "cybersecurity", "iss"],
    },
    "software engineering": {
        "course_code": "13477",
        "official_title": "Software Engineering",
        "credit_hours": 3,
        "prerequisites": ["Database Systems"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Covers software development life cycle, requirements engineering, software design, "
            "architecture patterns, testing, project management, and agile methodologies."
        ),
        "aliases": ["se"],
    },
    "natural language processing": {
        "course_code": "14351",
        "official_title": "Natural Language Processing",
        "credit_hours": 3,
        "prerequisites": ["Artificial Intelligence"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Covers core NLP concepts, linguistic and computational properties of natural "
            "language, and applications such as question answering, summarization, dialogue, "
            "and machine translation."
        ),
        "aliases": ["nlp"],
    },
    "computer architecture for machine learning": {
        "course_code": "14350",
        "official_title": "Computer Architecture for Machine Learning",
        "credit_hours": 3,
        "prerequisites": ["Artificial Intelligence"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Foundations of machine learning, implementation of algorithms and applications, including supervised "
            "learning, unsupervised learning, deep learning, reinforcement learning, and computer architectures "
            "for efficient execution such as CPU, GPU, and TensorFlow basics."
        ),
        "aliases": ["caml", "ml architecture"],
    },
    "business intelligence": {
        "course_code": "14466",
        "official_title": "Business Intelligence",
        "credit_hours": 3,
        "prerequisites": ["Data Visualization"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Introduces business intelligence concepts, analytics, enterprise data warehousing, "
            "decision support, and case studies in BI applications."
        ),
        "aliases": ["bi"],
    },
    "pattern recognition": {
        "course_code": "14452",
        "official_title": "Pattern Recognition",
        "credit_hours": 3,
        "prerequisites": ["Artificial Intelligence"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Covers statistical and structural approaches to pattern recognition, feature extraction, "
            "classification, clustering, neural networks for recognition, and applications in image "
            "and speech recognition."
        ),
        "aliases": [],
    },
    "cloud computing": {
        "course_code": "14467",
        "official_title": "Cloud Computing",
        "credit_hours": 3,
        "prerequisites": ["Operating Systems"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 1,
        "description": (
            "Introduction to cloud computing, data centers, virtualization, cloud storage, programming models, "
            "service models, design and management of data centers, data distribution, durability, consistency, and redundancy."
        ),
        "aliases": [],
    },
    "computer vision": {
        "course_code": "14458",
        "official_title": "Computer Vision",
        "credit_hours": 3,
        "prerequisites": ["Artificial Intelligence"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 2,
        "description": (
            "Fundamentals of image formation, camera imaging geometry, feature detection and matching, "
            "stereo algorithms, motion estimation and tracking, image classification with neural networks, "
            "and object detection and tracking."
        ),
        "aliases": ["cv"],
    },
    "information retrieval": {
        "course_code": "14457",
        "official_title": "Information Retrieval",
        "credit_hours": 3,
        "prerequisites": ["Database Systems"],
        "concurrent": [],
        "recommended_year": 4,
        "recommended_semester": 2,
        "description": (
            "Introduces the principles and techniques of information retrieval systems, "
            "including text processing, indexing, query processing, ranking models, "
            "vector space model, probabilistic retrieval models, evaluation metrics, "
            "and applications such as search engines and document retrieval."
        ),
        "aliases": ["ir"],
    },
}


# =========================================================
# HELPERS
# =========================================================

def safe_slug(name: str) -> str:
    return re.sub(r"[^a-zA-Z0-9_]+", "_", str(name).strip().lower()).strip("_")


def collect_course_folders(courses_root: Path) -> list[Path]:
    if not courses_root.exists():
        raise FileNotFoundError(f"Courses folder not found: {courses_root}")
    return sorted([p for p in courses_root.iterdir() if p.is_dir()])


def should_ignore_path(path: Path) -> bool:
    return any(part in IGNORED_DIR_NAMES for part in path.parts)


def normalize_title(name: str) -> str:
    cleaned = re.sub(r"[_\-]+", " ", str(name).lower()).strip()
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned


def resolve_course_metadata(course_name: str) -> dict:
    normalized = normalize_title(course_name)

    if normalized in COURSE_METADATA_MAP:
        return COURSE_METADATA_MAP[normalized]

    for meta in COURSE_METADATA_MAP.values():
        aliases = [normalize_title(a) for a in meta.get("aliases", [])]
        if normalized in aliases:
            return meta

    return {
        "course_code": "",
        "official_title": course_name,
        "credit_hours": None,
        "prerequisites": [],
        "concurrent": [],
        "recommended_year": None,
        "recommended_semester": None,
        "description": "",
        "aliases": [],
    }


def collect_supported_sources(course_dir: Path) -> list[Path]:
    files = []

    for path in course_dir.rglob("*"):
        if should_ignore_path(path):
            continue
        if not path.is_file():
            continue
        if path.name.startswith("._"):
            continue

        suffix = path.suffix.lower()
        if suffix in SUPPORTED_DIRECT_FILE_EXTENSIONS or suffix == ".zip":
            files.append(path)

    return sorted(files)


def normalize_line(line: str) -> str:
    return re.sub(r"\s+", " ", line.strip()).lower()


def extract_json_block(text: str):
    """
    Extract and parse the first valid JSON object or array from `text`.
    Uses raw_decode so nested structures (concepts with subtopics arrays)
    are parsed correctly — the old greedy regex approach failed on them.
    """
    decoder = json.JSONDecoder()
    for start in range(len(text)):
        if text[start] not in ('{', '['):
            continue
        try:
            obj, _ = decoder.raw_decode(text, start)
            return obj
        except json.JSONDecodeError:
            continue
    return None


# =========================================================
# FILE HASH CACHE  (OPT: skip unchanged files on re-runs)
# =========================================================

def _file_md5(path: Path) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _hash_cache_path(course_dir: Path) -> Path:
    return course_dir / "outputs" / ".source_hashes.json"


def _load_hash_cache(course_dir: Path) -> dict[str, str]:
    p = _hash_cache_path(course_dir)
    if p.exists():
        try:
            with open(p) as f:
                return json.load(f)
        except Exception:
            return {}
    return {}


def _save_hash_cache(course_dir: Path, cache: dict[str, str]) -> None:
    p = _hash_cache_path(course_dir)
    p.parent.mkdir(parents=True, exist_ok=True)
    with open(p, "w") as f:
        json.dump(cache, f, indent=2)


def _files_changed_since_last_run(
    course_dir: Path,
    files: list[Path],
) -> tuple[list[Path], dict[str, str]]:
    """
    Return (changed_files, full_new_cache).
    If FORCE_REGENERATE is True, all files are treated as changed.
    """
    if FORCE_REGENERATE:
        new_cache = {str(f.relative_to(course_dir)): _file_md5(f) for f in files}
        return files, new_cache

    old_cache = _load_hash_cache(course_dir)
    new_cache: dict[str, str] = {}
    changed: list[Path] = []

    for f in files:
        rel    = str(f.relative_to(course_dir))
        digest = _file_md5(f)
        new_cache[rel] = digest
        if old_cache.get(rel) != digest:
            changed.append(f)

    return changed, new_cache


# =========================================================
# LLM HELPERS
# =========================================================

def ask_llm(
    prompt: str,
    max_tokens: int = MAX_OUTPUT_TOKENS,
    *,
    debug_label: str = "llm",
) -> str:
    requested_max_tokens = max_tokens
    request_kwargs = {
        "model": MODEL_NAME,
        "messages": [
            {
                "role": "system",
                "content": "You are a precise academic assistant. Return clean structured output only.",
            },
            {
                "role": "user",
                "content": prompt,
            },
        ],
        "max_completion_tokens": max_tokens,
    }

    if MODEL_NAME.startswith("gpt-5"):
        request_kwargs["reasoning_effort"] = TEXT_LLM_REASONING_EFFORT
        request_kwargs["verbosity"] = "low"
        request_kwargs["max_completion_tokens"] = max(
            max_tokens,
            MIN_TEXT_COMPLETION_BUDGET,
        )

    debug_semantic(
        f"LLM request start label={debug_label!r} model={MODEL_NAME} "
        f"prompt_chars={len(prompt)} requested_max_tokens={requested_max_tokens} "
        f"actual_max_completion_tokens={request_kwargs['max_completion_tokens']} "
        f"reasoning_effort={request_kwargs.get('reasoning_effort')!r}"
    )
    response = _call_openai_with_retries(
        lambda: get_client().chat.completions.create(**request_kwargs)
    )
    choice = response.choices[0] if getattr(response, "choices", None) else None
    message = getattr(choice, "message", None) if choice else None
    content = getattr(message, "content", "") if message else ""
    finish_reason = getattr(choice, "finish_reason", None) if choice else None
    usage = getattr(response, "usage", None)
    usage_preview = usage.model_dump() if hasattr(usage, "model_dump") else usage
    if content is None:
        content = ""
    text = str(content).strip()
    debug_semantic(
        f"LLM response label={debug_label!r} finish_reason={finish_reason!r} "
        f"content_chars={len(text)} usage={usage_preview}"
    )
    debug_semantic(f"LLM raw preview label={debug_label!r}: {debug_preview(text)!r}")
    if not text:
        debug_semantic(
            f"EMPTY LLM CONTENT label={debug_label!r}; "
            f"response_preview={debug_preview(response)!r}"
        )
    return text


def optimize_image_for_vision(image_bytes: bytes) -> tuple[bytes, int, int]:
    image = Image.open(io.BytesIO(image_bytes))
    image.load()
    if image.mode not in ("RGB", "L"):
        image = image.convert("RGB")
    elif image.mode == "L":
        image = image.convert("RGB")

    if image.width > MAX_VISION_IMAGE_WIDTH:
        ratio = MAX_VISION_IMAGE_WIDTH / image.width
        new_size = (MAX_VISION_IMAGE_WIDTH, max(1, int(image.height * ratio)))
        image = image.resize(new_size, Image.Resampling.LANCZOS)

    out = io.BytesIO()
    image.save(out, format="JPEG", quality=VISION_JPEG_QUALITY, optimize=True)
    return out.getvalue(), image.width, image.height


def estimate_image_token_usage(width: int, height: int) -> int:
    # Rough high-detail image estimate based on 512px tiles.
    tiles_w = max(1, (width + 511) // 512)
    tiles_h = max(1, (height + 511) // 512)
    return 85 + (170 * tiles_w * tiles_h)


def encode_image_bytes_to_data_url(image_bytes: bytes, mime_type: str = "image/jpeg") -> str:
    encoded = base64.b64encode(image_bytes).decode("utf-8")
    return f"data:{mime_type};base64,{encoded}"


def should_run_vision_on_pdf_page(page, extracted_text: str) -> bool:
    if VISION_QUOTA_EXHAUSTED or _openai_quota_exhausted.is_set():
        _increment_vision_stat("vision_pages_skipped")
        _increment_vision_stat("vision_pages_skipped_quota")
        return False

    if _vision_budget_exhausted.is_set():
        _increment_vision_stat("vision_pages_skipped")
        return False

    if not ENABLE_VISION or not RUN_VISION_EXTRACTION:
        _increment_vision_stat("vision_pages_skipped")
        return False

    text_len      = len((extracted_text or "").strip())
    has_images    = bool(page.get_images(full=True))

    if text_len >= MIN_TEXT_LENGTH_FOR_VISION_SKIP:
        _increment_vision_stat("vision_pages_skipped")
        return False

    # Vision is a scarce OCR recovery tool. Vector drawing/path counts are
    # intentionally ignored because slide templates over-trigger them.
    return (
        _pdf_page_likely_scanned_or_handwritten(extracted_text)
        or (has_images and text_len < IMAGE_TEXT_LENGTH_FOR_VISION)
    )


def _pdf_page_image_count(page) -> int:
    try:
        return len(page.get_images(full=True))
    except Exception:
        return 0


def _pdf_page_likely_scanned_or_handwritten(extracted_text: str) -> bool:
    return len((extracted_text or "").strip()) < MIN_TEXT_LENGTH_FOR_VISION_SKIP


def _pdf_page_is_vision_candidate(page, extracted_text: str) -> bool:
    if (
        VISION_QUOTA_EXHAUSTED
        or _openai_quota_exhausted.is_set()
        or _vision_budget_exhausted.is_set()
        or not ENABLE_VISION
        or not RUN_VISION_EXTRACTION
    ):
        return False

    text_len = len((extracted_text or "").strip())
    has_images = _pdf_page_image_count(page) > 0

    if text_len >= MIN_TEXT_LENGTH_FOR_VISION_SKIP:
        return False

    return (
        _pdf_page_likely_scanned_or_handwritten(extracted_text)
        or (has_images and text_len < IMAGE_TEXT_LENGTH_FOR_VISION)
    )


def should_run_vision_on_slide(slide_text: str, has_visual_shape: bool) -> bool:
    if VISION_QUOTA_EXHAUSTED or _openai_quota_exhausted.is_set():
        _increment_vision_stat("vision_pages_skipped")
        _increment_vision_stat("vision_pages_skipped_quota")
        return False

    if _vision_budget_exhausted.is_set():
        _increment_vision_stat("vision_pages_skipped")
        return False

    if not ENABLE_VISION or not RUN_VISION_EXTRACTION:
        _increment_vision_stat("vision_pages_skipped")
        return False

    text_len = len((slide_text or "").strip())
    if text_len >= MIN_TEXT_LENGTH_FOR_VISION_SKIP:
        _increment_vision_stat("vision_pages_skipped")
        return False

    should_run = (
        text_len < MIN_TEXT_LENGTH_FOR_VISION_SKIP
        or (has_visual_shape and text_len < IMAGE_TEXT_LENGTH_FOR_VISION)
    )
    if not should_run:
        _increment_vision_stat("vision_pages_skipped")
    return should_run


def ask_vision_on_image_bytes(image_bytes: bytes, prompt: str, *, page_label: str = "") -> str:
    """Uses the Responses API with image input."""
    if VISION_QUOTA_EXHAUSTED or _openai_quota_exhausted.is_set():
        _increment_vision_stat("vision_pages_skipped")
        _increment_vision_stat("vision_pages_skipped_quota")
        return ""

    if _vision_budget_exhausted.is_set():
        _increment_vision_stat("vision_pages_skipped")
        return ""

    optimized_bytes, width, height = optimize_image_for_vision(image_bytes)
    estimated_tokens = estimate_image_token_usage(width, height)
    label = f" {page_label}" if page_label else ""
    print(
        f"[VISION TOKEN ESTIMATE]{label} "
        f"image={width}x{height} bytes={len(optimized_bytes)} "
        f"expected_tokens~{estimated_tokens}"
    )
    _increment_vision_stat("vision_estimated_tokens", estimated_tokens)
    image_url = encode_image_bytes_to_data_url(optimized_bytes, mime_type="image/jpeg")

    response = _call_openai_with_retries(
        lambda: get_client().responses.create(
            model=VISION_MODEL_NAME,
            input=[
                {
                    "role": "user",
                    "content": [
                        {"type": "input_text", "text": prompt},
                        {
                            "type": "input_image",
                            "image_url": image_url,
                            "detail": "high",
                        },
                    ],
                }
            ],
            max_output_tokens=MAX_VISION_OUTPUT_TOKENS,
        ),
        is_vision=True,
    )

    if response is None:
        return ""

    text = getattr(response, "output_text", None)
    if text:
        return text.strip()

    try:
        parts = []
        for item in response.output:
            if getattr(item, "type", None) == "message":
                for c in getattr(item, "content", []):
                    if getattr(c, "type", None) == "output_text":
                        parts.append(c.text)
        return "\n".join(parts).strip()
    except Exception:
        return ""


def build_visual_prompt(file_type: str, page_number: int,
                        is_handwritten: bool = False,
                        is_vector_rich: bool = False) -> str:
    return f"""
Extract readable academic text from this {file_type} page {page_number}. Briefly summarize diagrams or charts. Describe formulas shortly. Ignore decoration. Return concise plain text only.
""".strip()


def render_pdf_page_to_png_bytes(page, zoom: float = PDF_RENDER_ZOOM) -> bytes:
    matrix = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    return pix.tobytes("png")


def _extract_pptx_chart_as_text(shape) -> str:
    """
    Extract a PowerPoint native chart (Insert → Chart) as structured text.
    python-pptx can read chart type and series data directly from the XML.
    Returns a plain-text description that is meaningful for embeddings.
    """
    try:
        chart = shape.chart
        type_name = str(chart.chart_type).split(".")[-1].replace("_", " ").title()
        lines = [f"[Chart Type: {type_name}]"]

        for series in chart.series:
            name = series.name or "Series"
            try:
                vals = [str(round(v, 4)) if isinstance(v, float) else str(v)
                        for v in series.values if v is not None]
                if vals:
                    lines.append(f"  Series '{name}': {', '.join(vals[:15])}")
            except Exception:
                lines.append(f"  Series '{name}': (values unavailable)")

        try:
            cats = list(chart.plots[0].series[0].data_labels)
            cat_texts = [str(c) for c in cats if c][:10]
            if cat_texts:
                lines.append(f"  Categories: {', '.join(cat_texts)}")
        except Exception:
            pass

        try:
            title = chart.chart_title.text_frame.text.strip()
            if title:
                lines.insert(1, f"  Title: {title}")
        except Exception:
            pass

        return "\n".join(lines)
    except Exception as e:
        return f"[Chart: data extraction failed — {e}]"


def _extract_smartart_or_group_as_text(shape) -> str:
    """
    Extract text from SmartArt and group shapes, preserving spatial order.
    """
    try:
        child_texts = []

        if hasattr(shape, "shapes"):
            children = list(shape.shapes)
            children.sort(key=lambda s: (
                getattr(s, "top", 0) or 0,
                getattr(s, "left", 0) or 0,
            ))
            for child in children:
                if hasattr(child, "text") and child.text.strip():
                    child_texts.append(child.text.strip())
                if hasattr(child, "shapes"):
                    for grandchild in child.shapes:
                        if hasattr(grandchild, "text") and grandchild.text.strip():
                            child_texts.append(grandchild.text.strip())
        elif hasattr(shape, "text") and shape.text.strip():
            child_texts.append(shape.text.strip())

        if not child_texts:
            return ""

        try:
            xml = shape._element.xml
            label = "[SmartArt/Diagram]" if ("dgm:" in xml or "dsp:" in xml) else "[Group Shape]"
        except Exception:
            label = "[Group/SmartArt]"

        return f"{label}\n" + " → ".join(child_texts)
    except Exception as e:
        return f"[Group/SmartArt: extraction failed — {e}]"


def extract_picture_texts_from_slide(
    slide,
    run_vision: bool,
    *,
    failure_context: Optional[dict] = None,
) -> tuple[list[str], bool]:
    """
    Extract visual content from a PowerPoint slide.
    """
    visual_blocks = []
    has_picture = False

    for shape in slide.shapes:

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_picture = True
            if not run_vision:
                continue
            try:
                image_bytes = shape.image.blob
                prompt = """
Extract readable academic text from this slide image. Briefly summarize diagrams or charts. Describe formulas shortly. Ignore decoration. Return concise plain text only.
""".strip()
                page_label = ""
                if failure_context:
                    page_label = (
                        f"{failure_context.get('relative_path', '')} "
                        f"page={failure_context.get('page', '')}"
                    ).strip()
                visual_text = ask_vision_on_image_bytes(
                    image_bytes,
                    prompt,
                    page_label=page_label,
                )
                if visual_text:
                    _increment_vision_stat("ocr_recoveries")
                    visual_blocks.append(visual_text[:MAX_VISION_CHARS_PER_PAGE])
            except Exception as e:
                print(f"  Vision failed on slide image: {e}")
                if failure_context:
                    record_failed_vision_page({
                        **failure_context,
                        "error": str(e),
                    })

        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            has_picture = True
            chart_text = _extract_pptx_chart_as_text(shape)
            if chart_text:
                visual_blocks.append(chart_text)

        elif shape.shape_type in (
            MSO_SHAPE_TYPE.GROUP,
            MSO_SHAPE_TYPE.CANVAS,
        ):
            has_picture = True
            group_text = _extract_smartart_or_group_as_text(shape)
            if group_text:
                visual_blocks.append(group_text)

    return visual_blocks, has_picture


def combine_page_text(text: str, visual_text: str) -> str:
    text = (text or "").strip()
    visual_text = (visual_text or "").strip()

    if text and visual_text:
        return f"{text}\n\n[Visual content]\n{visual_text}".strip()
    if text:
        return text
    if visual_text:
        return f"[Visual content]\n{visual_text}".strip()
    return ""


# =========================================================
# FILE READING
# =========================================================

def build_doc_record(file_name: str, relative_path: str, file_type: str, pages: list[dict]) -> dict:
    return {
        "file_name": file_name,
        "relative_path": relative_path,
        "file_type": file_type,
        "num_pages": len(pages),
        "pages": pages,
        "full_text": "\n\n".join([p.get("combined_text", p.get("text", "")) for p in pages]).strip(),
    }


# =========================================================
# OPT: PARALLEL PAGE PROCESSING — PDF
# Each page's vision call is independent IO. ThreadPoolExecutor gives
# near-linear speedup: 20 pages × 3s sequential = 60s → ~10s parallel.
# =========================================================

def _process_pdf_page(args: tuple) -> dict:
    """Worker: process a single PDF page in a thread."""
    page, page_index, context = args
    page_num = page_index + 1
    text = page.get_text("text").strip()
    visual_text = ""
    allowed_vision_pages = context.get("allowed_vision_pages", set())

    if page_num in allowed_vision_pages and should_run_vision_on_pdf_page(page, text):
        try:
            image_bytes    = render_pdf_page_to_png_bytes(page)
            prompt         = build_visual_prompt(
                file_type="pdf", page_number=page_num,
            )
            visual_text = ask_vision_on_image_bytes(
                image_bytes,
                prompt,
                page_label=f"{context.get('relative_path', '')} page={page_num}",
            )
            visual_text = visual_text[:MAX_VISION_CHARS_PER_PAGE].strip()
            if visual_text:
                _increment_vision_stat("ocr_recoveries")
        except Exception as e:
            print(f"  Vision failed on PDF page {page_num}: {e}")
            record_failed_vision_page({
                **context,
                "page": page_num,
                "error": str(e),
            })
    else:
        _increment_vision_stat("vision_pages_skipped")

    if not visual_text:
        _increment_vision_stat("text_only_pages_processed")

    return {
        "page":          page_num,
        "text":          text,
        "visual_text":   visual_text,
        "combined_text": combine_page_text(text, visual_text),
    }


def _extract_pdf_pages_parallel(
    doc: fitz.Document,
    *,
    course_dir: Path,
    relative_path: str,
    file_name: str,
) -> list[dict]:
    """Process all PDF pages concurrently."""
    pages_data: list[dict] = [None] * len(doc)
    vision_candidates = []
    for i in range(len(doc)):
        page = doc[i]
        text = page.get_text("text").strip()
        if _pdf_page_is_vision_candidate(page, text):
            vision_candidates.append({
                "page": i + 1,
                "text_len": len(text),
                "image_count": _pdf_page_image_count(page),
            })

    vision_candidates.sort(key=lambda item: (item["text_len"], -item["image_count"], item["page"]))
    allowed_vision_pages = {
        item["page"]
        for item in vision_candidates[:MAX_VISION_PAGES_PER_DOCUMENT]
    }
    if len(vision_candidates) > MAX_VISION_PAGES_PER_DOCUMENT:
        print(
            f"  Vision page cap for {relative_path}: "
            f"selected {len(allowed_vision_pages)}/{len(vision_candidates)} candidate page(s)"
        )

    context = {
        "course_dir": str(course_dir),
        "course_name": course_dir.name,
        "file_name": file_name,
        "relative_path": relative_path,
        "file_type": "pdf",
        "allowed_vision_pages": allowed_vision_pages,
    }

    with ThreadPoolExecutor(max_workers=MAX_PAGE_WORKERS) as pool:
        futures = {
            pool.submit(_process_pdf_page, (doc[i], i, context)): i
            for i in range(len(doc))
        }
        for future in as_completed(futures):
            i = futures[future]
            try:
                pages_data[i] = future.result()
            except Exception as e:
                text = doc[i].get_text("text").strip()
                pages_data[i] = {
                    "page": i + 1, "text": text,
                    "visual_text": "", "combined_text": text,
                }
                print(f"  Page {i+1} failed, using text only: {e}")

    return pages_data


def extract_pdf_text_from_path(pdf_path: Path, course_dir: Path, relative_path: str) -> dict:
    doc = fitz.open(pdf_path)
    pages = _extract_pdf_pages_parallel(
        doc,
        course_dir=course_dir,
        relative_path=relative_path,
        file_name=pdf_path.name,
    )
    doc.close()
    return build_doc_record(pdf_path.name, relative_path, "pdf", pages)


def extract_pdf_text_from_bytes(file_name: str, course_dir: Path, relative_path: str, file_bytes: bytes) -> dict:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = _extract_pdf_pages_parallel(
        doc,
        course_dir=course_dir,
        relative_path=relative_path,
        file_name=file_name,
    )
    doc.close()
    return build_doc_record(file_name, relative_path, "pdf", pages)


# =========================================================
# OPT: PARALLEL SLIDE PROCESSING — PPTX
# Same principle as PDF: each slide's vision call is independent IO.
# =========================================================

def _extract_pptx_speaker_notes(slide) -> str:
    try:
        notes_slide = slide.notes_slide
        notes_tf    = notes_slide.notes_text_frame
        notes_text  = notes_tf.text.strip()
        if notes_text:
            return f"[Speaker Notes]\n{notes_text}"
    except Exception:
        pass
    return ""


def _pptx_slide_has_visual_shapes(slide) -> bool:
    for shape in slide.shapes:
        if shape.shape_type in (
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.GROUP,
            MSO_SHAPE_TYPE.CANVAS,
        ):
            return True
    return False


def _pptx_slide_has_raster_picture(slide) -> bool:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    return False


def _extract_pptx_slide_text(slide) -> str:
    slide_text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            slide_text_parts.append(shape.text.strip())

    text = "\n".join(slide_text_parts).strip()

    notes_text = _extract_pptx_speaker_notes(slide)
    if notes_text:
        text = f"{text}\n\n{notes_text}".strip()

    return text


def _pptx_slide_is_vision_candidate(slide, slide_text: str) -> bool:
    if (
        _vision_budget_exhausted.is_set()
        or not ENABLE_VISION
        or not RUN_VISION_EXTRACTION
    ):
        return False

    text_len = len((slide_text or "").strip())
    has_picture = _pptx_slide_has_raster_picture(slide)

    if text_len >= MIN_TEXT_LENGTH_FOR_VISION_SKIP:
        return False

    return (
        text_len < MIN_TEXT_LENGTH_FOR_VISION_SKIP
        or (has_picture and text_len < IMAGE_TEXT_LENGTH_FOR_VISION)
    )


def _process_pptx_slide_worker(args: tuple) -> dict:
    """Worker: process a single PPTX slide in a thread."""
    slide, slide_number, context = args

    text = _extract_pptx_slide_text(slide)

    has_visual = _pptx_slide_has_visual_shapes(slide)
    allowed_vision_pages = context.get("allowed_vision_pages", set())
    if slide_number in allowed_vision_pages:
        run_vision = should_run_vision_on_slide(text, has_visual)
    else:
        _increment_vision_stat("vision_pages_skipped")
        run_vision = False

    visual_blocks, _ = extract_picture_texts_from_slide(
        slide,
        run_vision=run_vision,
        failure_context={**context, "page": slide_number},
    )
    visual_text = "\n\n".join([b for b in visual_blocks if b.strip()]).strip()
    if not visual_text:
        _increment_vision_stat("text_only_pages_processed")
    combined_text = combine_page_text(text, visual_text)

    return {
        "page":          slide_number,
        "text":          text,
        "visual_text":   visual_text,
        "combined_text": combined_text,
    }


def _extract_pptx_slides_parallel(
    prs: Presentation,
    *,
    course_dir: Path,
    relative_path: str,
    file_name: str,
) -> list[dict]:
    """Process all PPTX slides concurrently."""
    slides_data: list[dict] = [None] * len(prs.slides)
    vision_candidates = []
    for i, slide in enumerate(prs.slides):
        text = _extract_pptx_slide_text(slide)
        if _pptx_slide_is_vision_candidate(slide, text):
            vision_candidates.append({
                "page": i + 1,
                "text_len": len(text),
                "image_count": 1 if _pptx_slide_has_raster_picture(slide) else 0,
            })

    vision_candidates.sort(key=lambda item: (item["text_len"], -item["image_count"], item["page"]))
    allowed_vision_pages = {
        item["page"]
        for item in vision_candidates[:MAX_VISION_PAGES_PER_DOCUMENT]
    }
    if len(vision_candidates) > MAX_VISION_PAGES_PER_DOCUMENT:
        print(
            f"  Vision page cap for {relative_path}: "
            f"selected {len(allowed_vision_pages)}/{len(vision_candidates)} candidate slide(s)"
        )

    context = {
        "course_dir": str(course_dir),
        "course_name": course_dir.name,
        "file_name": file_name,
        "relative_path": relative_path,
        "file_type": "pptx",
        "allowed_vision_pages": allowed_vision_pages,
    }

    with ThreadPoolExecutor(max_workers=MAX_PAGE_WORKERS) as pool:
        futures = {
            pool.submit(_process_pptx_slide_worker, (slide, i + 1, context)): i
            for i, slide in enumerate(prs.slides)
        }
        for future in as_completed(futures):
            i = futures[future]
            try:
                slides_data[i] = future.result()
            except Exception as e:
                slides_data[i] = {
                    "page": i + 1, "text": "", "visual_text": "", "combined_text": "",
                }
                print(f"  Slide {i+1} failed: {e}")

    return slides_data


def extract_pptx_text_from_path(pptx_path: Path, course_dir: Path, relative_path: str) -> dict:
    prs = Presentation(str(pptx_path))
    return build_doc_record(
        pptx_path.name, relative_path, "pptx",
        _extract_pptx_slides_parallel(
            prs,
            course_dir=course_dir,
            relative_path=relative_path,
            file_name=pptx_path.name,
        ),
    )


def extract_pptx_text_from_bytes(file_name: str, course_dir: Path, relative_path: str, file_bytes: bytes) -> dict:
    prs = Presentation(io.BytesIO(file_bytes))
    return build_doc_record(
        file_name, relative_path, "pptx",
        _extract_pptx_slides_parallel(
            prs,
            course_dir=course_dir,
            relative_path=relative_path,
            file_name=file_name,
        ),
    )


def extract_code_text_from_path(file_path: Path, relative_path: str) -> dict:
    text = file_path.read_text(encoding="utf-8", errors="ignore").strip()
    return build_doc_record(
        file_name=file_path.name,
        relative_path=relative_path,
        file_type=file_path.suffix.lower().lstrip("."),
        pages=[{"page": 1, "text": text, "visual_text": "", "combined_text": text}],
    )


def extract_code_text_from_bytes(file_name: str, relative_path: str, file_type: str, file_bytes: bytes) -> dict:
    text = file_bytes.decode("utf-8", errors="ignore").strip()
    return build_doc_record(
        file_name=file_name,
        relative_path=relative_path,
        file_type=file_type,
        pages=[{"page": 1, "text": text, "visual_text": "", "combined_text": text}],
    )


def extract_regular_file(file_path: Path, course_dir: Path) -> list[dict]:
    suffix        = file_path.suffix.lower()
    relative_path = str(file_path.relative_to(course_dir))

    if suffix == ".pdf":
        return [extract_pdf_text_from_path(file_path, course_dir, relative_path)]
    if suffix == ".pptx":
        return [extract_pptx_text_from_path(file_path, course_dir, relative_path)]

    return [extract_code_text_from_path(file_path, relative_path)]


def extract_zip_contents(zip_path: Path, course_dir: Path) -> list[dict]:
    docs         = []
    zip_relative = str(zip_path.relative_to(course_dir))

    with zipfile.ZipFile(zip_path, "r") as zf:
        for member_name in zf.namelist():
            if member_name.endswith("/"):
                continue

            member_path = Path(member_name)

            if member_path.name.startswith("._"):
                continue
            if any(part in IGNORED_DIR_NAMES for part in member_path.parts):
                continue

            suffix = member_path.suffix.lower()
            if suffix not in SUPPORTED_ZIP_MEMBER_EXTENSIONS:
                continue

            file_bytes            = zf.read(member_name)
            virtual_relative_path = f"{zip_relative}::{member_name}"

            try:
                if suffix == ".pdf":
                    docs.append(extract_pdf_text_from_bytes(
                        file_name=member_path.name,
                        course_dir=course_dir,
                        relative_path=virtual_relative_path,
                        file_bytes=file_bytes,
                    ))
                elif suffix == ".pptx":
                    docs.append(extract_pptx_text_from_bytes(
                        file_name=member_path.name,
                        course_dir=course_dir,
                        relative_path=virtual_relative_path,
                        file_bytes=file_bytes,
                    ))
                else:
                    docs.append(extract_code_text_from_bytes(
                        file_name=member_path.name,
                        relative_path=virtual_relative_path,
                        file_type=suffix.lstrip("."),
                        file_bytes=file_bytes,
                    ))
            except Exception as e:
                print(f"  Skipped zip member {virtual_relative_path}: {e}")

    return docs


def extract_source_docs(file_path: Path, course_dir: Path) -> list[dict]:
    suffix = file_path.suffix.lower()
    if suffix == ".zip":
        return extract_zip_contents(file_path, course_dir)
    return extract_regular_file(file_path, course_dir)


def extract_source_docs_with_timing(file_path: Path, course_dir: Path) -> list[dict]:
    start = time.time()
    try:
        return extract_source_docs(file_path, course_dir)
    finally:
        elapsed = time.time() - start
        try:
            file_label = str(file_path.relative_to(course_dir))
        except ValueError:
            file_label = str(file_path)
        print(f"[FILE TIME] {file_label} took {elapsed:.2f}s")


def load_all_docs(files: list[Path], course_dir: Path) -> list[dict]:
    docs = []
    _increment_vision_stat("total_files", len(files))

    with ThreadPoolExecutor(max_workers=MAX_FILE_WORKERS) as pool:
        futures = {pool.submit(extract_source_docs_with_timing, f, course_dir): f for f in files}
        for future in tqdm(
            as_completed(futures),
            total=len(futures),
            desc="Reading files",
            leave=False,
        ):
            f = futures[future]
            try:
                docs.extend(future.result())
            except Exception as e:
                print(f"  Skipped {f.name}: {e}")

    return docs


def _dedupe_failed_vision_entries(entries: list[dict]) -> list[dict]:
    seen = set()
    deduped = []

    for entry in entries:
        key = (
            entry.get("course_dir", ""),
            entry.get("relative_path", ""),
            entry.get("file_type", ""),
            entry.get("page"),
        )
        if key in seen:
            continue
        seen.add(key)
        deduped.append(entry)

    return deduped


def _source_bytes_for_relative_path(course_dir: Path, relative_path: str) -> tuple[str, bytes]:
    if "::" not in relative_path:
        source_path = course_dir / relative_path
        return source_path.name, source_path.read_bytes()

    zip_relative, member_name = relative_path.split("::", 1)
    zip_path = course_dir / zip_relative
    with zipfile.ZipFile(zip_path, "r") as zf:
        return Path(member_name).name, zf.read(member_name)


def _recover_pdf_page_visual_text(course_dir: Path, entry: dict) -> str:
    _, file_bytes = _source_bytes_for_relative_path(course_dir, entry["relative_path"])
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        page_number = int(entry["page"])
        page = doc[page_number - 1]
        text = page.get_text("text").strip()
        image_bytes = render_pdf_page_to_png_bytes(page)
        prompt = build_visual_prompt(
            file_type="pdf",
            page_number=page_number,
        )
        return ask_vision_on_image_bytes(
            image_bytes,
            prompt,
            page_label=f"{entry.get('relative_path', '')} page={page_number} recovery",
        )[:MAX_VISION_CHARS_PER_PAGE].strip()
    finally:
        doc.close()


def _recover_pptx_slide_visual_text(course_dir: Path, entry: dict) -> str:
    _, file_bytes = _source_bytes_for_relative_path(course_dir, entry["relative_path"])
    prs = Presentation(io.BytesIO(file_bytes))
    slide_number = int(entry["page"])
    slide = prs.slides[slide_number - 1]
    visual_blocks, _ = extract_picture_texts_from_slide(
        slide,
        run_vision=True,
        failure_context=None,
    )
    return "\n\n".join([b for b in visual_blocks if b.strip()])[:MAX_VISION_CHARS_PER_PAGE].strip()


def _apply_recovered_visual_text(docs: list[dict], entry: dict, visual_text: str) -> bool:
    if not visual_text:
        return False

    for doc in docs:
        if doc.get("relative_path") != entry.get("relative_path"):
            continue
        for page in doc.get("pages", []):
            if int(page.get("page", -1)) != int(entry.get("page", -2)):
                continue

            existing_visual = page.get("visual_text", "").strip()
            page["visual_text"] = (
                f"{existing_visual}\n\n{visual_text}".strip()
                if existing_visual and visual_text not in existing_visual
                else visual_text
            )
            page["combined_text"] = combine_page_text(page.get("text", ""), page["visual_text"])
            doc["full_text"] = "\n\n".join(
                [p.get("combined_text", p.get("text", "")) for p in doc.get("pages", [])]
            ).strip()
            return True

    return False


def run_failed_vision_recovery_pass(course_dir: Path, docs: list[dict]) -> None:
    """
    Re-run pages that exhausted all retries during parallel extraction. This is
    intentionally sequential, so recovered visual text can feed chunking, KG
    concept extraction, and vectorization for the current course.
    """
    with _failed_vision_lock:
        entries = [
            entry for entry in _load_failed_vision_pages()
            if entry.get("course_dir") == str(course_dir)
        ]
    entries = _dedupe_failed_vision_entries(entries)

    if not entries:
        return

    print(f"  Recovery pass: retrying {len(entries)} failed vision page(s) sequentially")
    recovered_keys = set()

    for entry in entries:
        key = (
            entry.get("course_dir", ""),
            entry.get("relative_path", ""),
            entry.get("file_type", ""),
            entry.get("page"),
        )
        try:
            if entry.get("file_type") == "pdf":
                visual_text = _recover_pdf_page_visual_text(course_dir, entry)
            elif entry.get("file_type") == "pptx":
                visual_text = _recover_pptx_slide_visual_text(course_dir, entry)
            else:
                visual_text = ""

            if _apply_recovered_visual_text(docs, entry, visual_text):
                _increment_vision_stat("ocr_recoveries")
                recovered_keys.add(key)
                print(
                    f"    recovered {entry.get('relative_path')} "
                    f"page/slide {entry.get('page')}"
                )
            else:
                print(
                    f"    recovery returned no visual text for {entry.get('relative_path')} "
                    f"page/slide {entry.get('page')}"
                )
        except Exception as e:
            print(
                f"    recovery failed for {entry.get('relative_path')} "
                f"page/slide {entry.get('page')}: {e}"
            )

    if recovered_keys:
        with _failed_vision_lock:
            unresolved = []
            for entry in _load_failed_vision_pages():
                key = (
                    entry.get("course_dir", ""),
                    entry.get("relative_path", ""),
                    entry.get("file_type", ""),
                    entry.get("page"),
                )
                if key not in recovered_keys:
                    unresolved.append(entry)
            _write_failed_vision_pages(unresolved)

    print_vision_counters("  Vision counters after recovery")


# =========================================================
# METADATA JSON WRITING
# =========================================================

def build_materials_index(course_dir: Path, docs_df: pd.DataFrame) -> dict:
    materials = []
    for i, (_, row) in enumerate(docs_df.iterrows(), start=1):
        materials.append({
            "material_id":    f"{safe_slug(course_dir.name)}_{i:03d}",
            "file_name":      row["file_name"],
            "relative_path":  row["relative_path"],
            "file_type":      row["file_type"],
            "chapter":        row["chapter"],
        })
    return {"course_folder": course_dir.name, "materials": materials}


def write_metadata_files(course_dir: Path, docs_df: pd.DataFrame):
    metadata_dir = course_dir / "metadata"
    metadata_dir.mkdir(parents=True, exist_ok=True)

    course_meta = resolve_course_metadata(course_dir.name)

    course_info = {
        "course_code":           course_meta["course_code"],
        "course_title":          course_meta["official_title"],
        "course_folder_name":    course_dir.name,
        "credit_hours":          course_meta["credit_hours"],
        "prerequisites":         course_meta["prerequisites"],
        "concurrent":            course_meta["concurrent"],
        "recommended_year":      course_meta["recommended_year"],
        "recommended_semester":  course_meta["recommended_semester"],
        "description":           course_meta["description"],
    }

    prerequisites = {
        "course_code":  course_meta["course_code"],
        "course_title": course_meta["official_title"],
        "prerequisites": course_meta["prerequisites"],
        "concurrent":   course_meta["concurrent"],
    }

    materials_index = build_materials_index(course_dir, docs_df)

    topic_schema = {
        "schema_name": "topic_object_v2",
        "fields": [
            "topic_name", "subtopics", "file_source",
            "relative_path", "chapter", "summary", "keywords",
        ],
    }

    files = {
        "course_info.json":     course_info,
        "prerequisites.json":   prerequisites,
        "materials_index.json": materials_index,
        "topic_schema.json":    topic_schema,
    }

    for filename, content in files.items():
        with open(metadata_dir / filename, "w", encoding="utf-8") as f:
            json.dump(content, f, indent=2, ensure_ascii=False)

    return {
        "course_info":     course_info,
        "prerequisites":   prerequisites,
        "materials_index": materials_index,
        "topic_schema":    topic_schema,
    }


# =========================================================
# CLEANING
# =========================================================

_BAD_EXACT_LINES: set[str] = {
    "copyright", "references", "reference", "lecture notes",
    "course instructor", "any question", "thank you",
    "princess sumaya university for technology", "king hussein school",
    "data science department", "amman", "p.o.box", "al-jubaiha",
    "www.psut.edu.jo", "info@psut.edu.jo", "psut", "dr.", "prof.",
    "ibrahim abu alhaol", "bushra alhijawi", "ahmad alzghoul",
    "mohammad azzeh", "ibrahim", "bushra", "alhijawi", "alzghoul",
}

_BAD_PATTERN_REGEXES: list[re.Pattern] = [
    re.compile(r"\bwww\b"),
    re.compile(r"\bfax\b"),
    re.compile(r"\b962\b"),
    re.compile(r"@psut\.edu\.jo"),
    re.compile(r"\.edu\.jo\b"),
]

def _is_bad_line(norm: str) -> bool:
    if norm in _BAD_EXACT_LINES:
        return True
    for pat in _BAD_PATTERN_REGEXES:
        if pat.search(norm):
            return True
    return False

bad_patterns: list[str] = []  # legacy alias, unused


def remove_repeated_page_lines(docs: list[dict], min_repeat_ratio: float = 0.35):
    all_lines  = []
    page_count = 0

    for doc in docs:
        if doc["file_type"] not in {"pdf", "pptx"}:
            continue
        for page in doc["pages"]:
            page_count += 1
            lines = [normalize_line(l) for l in page.get("text", "").splitlines() if l.strip()]
            all_lines.extend(set(lines))

    line_counts    = Counter(all_lines)
    repeated_lines = {
        line for line, count in line_counts.items()
        if page_count > 0 and count >= page_count * min_repeat_ratio
    }

    cleaned_docs = []

    for doc in docs:
        cleaned_pages = []

        for page in doc["pages"]:
            original_lines = page.get("text", "").splitlines()
            kept_lines     = []

            for line in original_lines:
                norm = normalize_line(line)

                if not norm:
                    continue
                if doc["file_type"] in {"pdf", "pptx"} and norm in repeated_lines:
                    continue
                if doc["file_type"] in {"pdf", "pptx"} and _is_bad_line(norm):
                    continue
                if doc["file_type"] in {"pdf", "pptx"} and re.fullmatch(r"\d{1,4}", norm):
                    continue
                if len(norm) <= 2:
                    continue
                if norm.isdigit():
                    continue

                kept_lines.append(line)

            cleaned_text  = "\n".join(kept_lines).strip()
            visual_text   = page.get("visual_text", "").strip()
            combined_text = combine_page_text(cleaned_text, visual_text)

            cleaned_pages.append({
                "page":          page["page"],
                "text":          cleaned_text,
                "visual_text":   visual_text,
                "combined_text": combined_text,
            })

        cleaned_docs.append({
            "file_name":     doc["file_name"],
            "relative_path": doc["relative_path"],
            "file_type":     doc["file_type"],
            "num_pages":     doc["num_pages"],
            "pages":         cleaned_pages,
            "full_text":     "\n\n".join([p.get("combined_text", "") for p in cleaned_pages]).strip(),
        })

    return cleaned_docs, repeated_lines


# =========================================================
# DATAFRAMES + CHUNKING
# =========================================================

def build_docs_df(docs: list[dict]) -> pd.DataFrame:
    df = pd.DataFrame(docs)
    df["chapter"] = df["file_name"].str.extract(r"(CH\d+|ch\d+)", expand=False)
    df["chapter"] = df["chapter"].fillna("").str.upper()
    return df


def chunk_text(text: str, chunk_size: int = 450, overlap: int = 80) -> list[str]:
    words  = text.split()
    chunks = []
    start  = 0

    while start < len(words):
        end = min(len(words), start + chunk_size)
        chunks.append(" ".join(words[start:end]))
        if end == len(words):
            break
        start = end - overlap

    return chunks


def build_chunks_df(course_name: str, docs_df: pd.DataFrame) -> pd.DataFrame:
    rows = []
    for _, row in docs_df.iterrows():
        for idx, chunk in enumerate(chunk_text(row["full_text"])):
            rows.append({
                "course_name":   course_name,
                "file_name":     row["file_name"],
                "relative_path": row["relative_path"],
                "file_type":     row["file_type"],
                "chapter":       row["chapter"],
                "chunk_id":      idx,
                "chunk_text":    chunk,
            })
    return pd.DataFrame(rows)


# =========================================================
# OPT: PARALLEL SUMMARIZATION
# Sequential: 43 files × 3s = 129s. All pure IO-wait.
# Parallel with 8 workers: ~18s (~7x speedup).
# =========================================================

def _summarize_one(args: tuple) -> tuple[str, str]:
    """Worker: summarize a single document. Returns (relative_path, summary)."""
    course_name, relative_path, full_text = args
    debug_semantic(
        f"Summary generation start course={course_name!r} file={relative_path!r} "
        f"full_text_chars={len(full_text or '')}"
    )
    # OPT: context trimmed 6000→3000 chars; tokens 900→600
    prompt = f"""
Write a detailed academic summary of this course file.

Requirements:
- 2 to 4 medium-length paragraphs
- not bullet points
- mention main concepts, important methods, and what the file teaches
- if the content contains mathematical formulas or statistical notation, briefly describe
  what the key formulas compute in plain English
- if the content describes diagrams, charts, or algorithms, mention their purpose
- keep it clear and study-oriented
- do not add information not found in the content

Course: {course_name}
File: {relative_path}

Content:
{full_text[:3000]}
"""
    summary = ask_llm(
        prompt,
        max_tokens=600,
        debug_label=f"summary:{course_name}:{relative_path}",
    )
    debug_semantic(
        f"Summary parsed course={course_name!r} file={relative_path!r} "
        f"summary_chars={len(summary)} preview={debug_preview(summary)!r}"
    )
    return relative_path, summary


def run_summarization(course_name: str, docs_df: pd.DataFrame) -> dict:
    """OPT: parallel summarization across all files in the course."""
    args_list = [
        (course_name, row["relative_path"], row["full_text"])
        for _, row in docs_df.iterrows()
    ]
    summaries: dict[str, str] = {}

    with ThreadPoolExecutor(max_workers=MAX_FILE_WORKERS) as pool:
        futures = {pool.submit(_summarize_one, args): args[1] for args in args_list}
        for future in tqdm(as_completed(futures), total=len(futures),
                           desc=f"Summarizing {course_name}", leave=False):
            try:
                rel_path, summary = future.result()
                summaries[rel_path] = summary
            except Exception as e:
                rel_path = futures[future]
                summaries[rel_path] = ""
                print(f"  Summarization failed for {rel_path}: {type(e).__name__}: {e!r}")

    nonempty = sum(bool(str(v).strip()) for v in summaries.values())
    debug_semantic(
        f"Summarization complete course={course_name!r} total={len(summaries)} "
        f"nonempty={nonempty} empty={len(summaries) - nonempty}"
    )
    return summaries


# =========================================================
# OPT: PARALLEL CONCEPT EXTRACTION
# Sequential: 43 files × 3.5s = 150s. All pure IO-wait.
# Parallel with 8 workers: ~22s (~7x speedup).
# =========================================================

def _extract_concepts_one(args: tuple) -> tuple[str, dict]:
    """Worker: extract concepts from a single document."""
    course_name, relative_path, chapter, full_text = args
    debug_semantic(
        f"Concept generation start course={course_name!r} file={relative_path!r} "
        f"chapter={chapter!r} full_text_chars={len(full_text or '')}"
    )
    # OPT: context trimmed 5000→2500 chars; tokens 900→500
    prompt = f"""
You are extracting structured academic concepts from university course content.

Course: {course_name}
File: {relative_path}

Return ONLY valid JSON in this exact format:
{{
  "chapter": "{chapter}",
  "topics": [
    {{
      "topic_name": "...",
      "subtopics": ["...", "..."],
      "keywords": ["...", "..."],
      "formula_descriptions": ["...", "..."],
      "visual_types": ["...", "..."]
    }}
  ]
}}

Rules:
- Return valid JSON only. No text before or after the JSON.
- Extract only concepts actually present in the content.
- topic_name: concise, academic, specific (e.g. "Maximum Likelihood Estimation", not "Statistics").
- subtopics: specific sub-concepts under the topic (e.g. "Log-likelihood derivation", "MLE for Gaussian").
- keywords: short technical terms (e.g. "likelihood", "gradient", "convergence").
- formula_descriptions: for EVERY formula, equation, or mathematical notation present,
  write a plain-English description of what it computes or expresses.
  Example: "E[X] = Σ xᵢP(xᵢ)" → "Expected value as a probability-weighted average of outcomes".
  This field may be an empty list [] if no formulas are present.
- visual_types: list the types of visual elements present on this page/slide.
  Use labels from: ["chart", "graph", "flowchart", "diagram", "table", "code", "pseudocode",
  "tree", "matrix", "formula", "handwriting", "screenshot", "none"].
  This field may be ["none"] if no visuals are present.

Content:
{full_text[:2500]}
"""
    raw    = ask_llm(
        prompt,
        max_tokens=500,
        debug_label=f"concepts:{course_name}:{relative_path}",
    )
    debug_semantic(
        f"Concept raw response course={course_name!r} file={relative_path!r} "
        f"raw_chars={len(raw)} preview={debug_preview(raw)!r}"
    )
    parsed = extract_json_block(raw)

    if parsed and isinstance(parsed, dict) and "topics" in parsed:
        topics = parsed.get("topics", [])
        topic_count = len(topics) if isinstance(topics, list) else 0
        subtopic_count = 0
        if isinstance(topics, list):
            for topic in topics:
                if isinstance(topic, dict) and isinstance(topic.get("subtopics"), list):
                    subtopic_count += len([s for s in topic["subtopics"] if str(s).strip()])
        debug_semantic(
            f"Concept parsed course={course_name!r} file={relative_path!r} "
            f"topics={topic_count} subtopics={subtopic_count} "
            f"parsed_preview={debug_preview(json.dumps(parsed, ensure_ascii=False))!r}"
        )
        return relative_path, parsed
    debug_semantic(
        f"Concept parse failed/empty course={course_name!r} file={relative_path!r}; "
        f"returning empty topics. raw_preview={debug_preview(raw)!r}"
    )
    return relative_path, {"chapter": chapter, "topics": []}


def run_concept_extraction(course_name: str, docs_df: pd.DataFrame) -> dict:
    """OPT: parallel concept extraction across all files in the course."""
    args_list = [
        (course_name, row["relative_path"], row["chapter"], row["full_text"])
        for _, row in docs_df.iterrows()
    ]
    concepts: dict[str, dict] = {}

    with ThreadPoolExecutor(max_workers=MAX_FILE_WORKERS) as pool:
        futures = {pool.submit(_extract_concepts_one, args): args[1] for args in args_list}
        for future in tqdm(as_completed(futures), total=len(futures),
                           desc=f"Concept extraction {course_name}", leave=False):
            try:
                rel_path, result = future.result()
                concepts[rel_path] = result
            except Exception as e:
                rel_path = futures[future]
                concepts[rel_path] = {"chapter": "", "topics": []}
                print(f"  Concept extraction failed for {rel_path}: {type(e).__name__}: {e!r}")

    topic_count, subtopic_count = count_concept_items(concepts)
    debug_semantic(
        f"Concept extraction complete course={course_name!r} files={len(concepts)} "
        f"topics={topic_count} subtopics={subtopic_count}"
    )
    return concepts


# =========================================================
# DATA PREP FOR SUMMARY / CONCEPT / METADATA VECTORIZATION
# =========================================================

def build_summaries_df(course_name: str, summaries: dict) -> pd.DataFrame:
    rows = []
    for relative_path, summary_text in summaries.items():
        rows.append({
            "course_name":   course_name,
            "relative_path": relative_path,
            "summary_text":  summary_text,
        })
    return pd.DataFrame(rows)


def flatten_concepts_for_chroma(course_name: str, concepts: dict) -> pd.DataFrame:
    source_topics, source_subtopics = count_concept_items(concepts)
    debug_semantic(
        f"Flatten concepts start course={course_name!r}: "
        f"source_files={len(concepts) if isinstance(concepts, dict) else 'n/a'} "
        f"source_topics={source_topics} source_subtopics={source_subtopics}"
    )
    rows = []

    for relative_path, content in concepts.items():
        chapter = content.get("chapter", "")

        for idx, topic in enumerate(content.get("topics", [])):
            topic_name           = topic.get("topic_name", "").strip()
            subtopics            = [str(s).strip() for s in topic.get("subtopics",            []) if str(s).strip()]
            keywords             = [str(k).strip() for k in topic.get("keywords",             []) if str(k).strip()]
            formula_descriptions = [str(f).strip() for f in topic.get("formula_descriptions", []) if str(f).strip()]
            visual_types         = [str(v).strip() for v in topic.get("visual_types",         []) if str(v).strip()]

            parts = [
                f"Course: {course_name}",
                f"Chapter: {chapter}",
                f"Topic: {topic_name}",
            ]
            if subtopics:
                parts.append(f"Subtopics: {', '.join(subtopics)}")
            if keywords:
                parts.append(f"Keywords: {', '.join(keywords)}")
            if formula_descriptions:
                parts.append(f"Formulas: {' | '.join(formula_descriptions)}")
            if visual_types and visual_types != ["none"]:
                parts.append(f"Visual content: {', '.join(visual_types)}")

            concept_text = "\n".join(parts).strip()

            rows.append({
                "course_name":      course_name,
                "relative_path":    relative_path,
                "chapter":          chapter,
                "topic_index":      idx,
                "topic_name":       topic_name,
                "concept_text":     concept_text,
                "has_formulas":     len(formula_descriptions) > 0,
                "has_visuals":      bool(visual_types and visual_types != ["none"]),
                "visual_types_str": ", ".join(visual_types) if visual_types else "none",
            })

    df = pd.DataFrame(rows)
    debug_semantic(f"Flatten concepts complete course={course_name!r}: rows={len(df)}")
    return df


def build_metadata_documents(course_name: str, metadata_objects: dict) -> pd.DataFrame:
    rows = []
    for doc_type, content in metadata_objects.items():
        rows.append({
            "course_name":   course_name,
            "doc_type":      doc_type,
            "document_text": json.dumps(content, ensure_ascii=False, indent=2),
        })
    return pd.DataFrame(rows)


# =========================================================
# CHROMADB / VECTORIZATION
# =========================================================

def get_chroma_client(chroma_dir: Path):
    chroma_dir.mkdir(parents=True, exist_ok=True)
    return chromadb.PersistentClient(path=str(chroma_dir))


def get_embedding_function():
    return embedding_functions.DefaultEmbeddingFunction()


def recreate_collection(chroma_client, collection_name: str, embedding_function):
    try:
        chroma_client.delete_collection(name=collection_name)
    except Exception:
        pass
    return chroma_client.get_or_create_collection(
        name=collection_name,
        embedding_function=embedding_function,
    )


def save_chunks_to_chroma(course_name: str, chunks_df: pd.DataFrame, chroma_dir: Path):
    chroma_client      = get_chroma_client(chroma_dir)
    embedding_function = get_embedding_function()
    collection_name    = f"{safe_slug(course_name)}{CHUNKS_COLLECTION_SUFFIX}"
    collection         = recreate_collection(chroma_client, collection_name, embedding_function)

    ids, documents, metadatas = [], [], []
    for _, row in chunks_df.iterrows():
        ids.append(f"{safe_slug(course_name)}_{safe_slug(row['relative_path'])}_{row['chunk_id']}")
        documents.append(row["chunk_text"])
        metadatas.append({
            "doc_type":      "chunk",
            "course_name":   str(row["course_name"]),
            "file_name":     str(row["file_name"]),
            "relative_path": str(row["relative_path"]),
            "file_type":     str(row["file_type"]),
            "chapter":       str(row["chapter"]),
            "chunk_id":      int(row["chunk_id"]),
        })

    collection.upsert(ids=ids, documents=documents, metadatas=metadatas)
    print(f"  Chunks vectorized and saved: {collection.count()}")


def save_summaries_to_chroma(course_name: str, summaries_df: pd.DataFrame, chroma_dir: Path):
    chroma_client      = get_chroma_client(chroma_dir)
    embedding_function = get_embedding_function()
    collection         = recreate_collection(chroma_client, f"{safe_slug(course_name)}_summaries", embedding_function)

    ids, documents, metadatas = [], [], []
    for idx, row in summaries_df.iterrows():
        ids.append(f"{safe_slug(course_name)}_summary_{idx}")
        documents.append(row["summary_text"])
        metadatas.append({
            "doc_type":      "summary",
            "course_name":   str(row["course_name"]),
            "relative_path": str(row["relative_path"]),
        })

    if ids:
        collection.upsert(ids=ids, documents=documents, metadatas=metadatas)
    print(f"  Summaries vectorized and saved: {collection.count()}")


def save_concepts_to_chroma(course_name: str, concepts_df: pd.DataFrame, chroma_dir: Path):
    debug_semantic(
        f"Save concepts to Chroma start course={course_name!r}: rows={len(concepts_df)} "
        f"columns={list(concepts_df.columns)}"
    )
    chroma_client      = get_chroma_client(chroma_dir)
    embedding_function = get_embedding_function()
    collection         = recreate_collection(chroma_client, f"{safe_slug(course_name)}_concepts", embedding_function)

    ids, documents, metadatas = [], [], []
    for idx, row in concepts_df.iterrows():
        ids.append(f"{safe_slug(course_name)}_concept_{idx}")
        documents.append(row["concept_text"])
        metadatas.append({
            "doc_type":      "concept",
            "course_name":   str(row["course_name"]),
            "relative_path": str(row["relative_path"]),
            "chapter":       str(row["chapter"]),
            "topic_name":    str(row["topic_name"]),
            "topic_index":   int(row["topic_index"]),
            "has_formulas":  bool(row.get("has_formulas", False)),
            "has_visuals":   bool(row.get("has_visuals", False)),
            "visual_types":  str(row.get("visual_types_str", "none")),
        })

    if ids:
        collection.upsert(ids=ids, documents=documents, metadatas=metadatas)
    else:
        debug_semantic(
            f"Save concepts to Chroma skipped upsert course={course_name!r}: no concept rows"
        )
    print(f"  Concepts vectorized and saved: {collection.count()}")


def save_metadata_to_chroma(course_name: str, metadata_df: pd.DataFrame, chroma_dir: Path):
    chroma_client      = get_chroma_client(chroma_dir)
    embedding_function = get_embedding_function()
    collection         = recreate_collection(chroma_client, f"{safe_slug(course_name)}_metadata", embedding_function)

    ids, documents, metadatas = [], [], []
    for idx, row in metadata_df.iterrows():
        ids.append(f"{safe_slug(course_name)}_metadata_{idx}")
        documents.append(row["document_text"])
        metadatas.append({
            "doc_type":      "metadata",
            "course_name":   str(row["course_name"]),
            "metadata_type": str(row["doc_type"]),
        })

    if ids:
        collection.upsert(ids=ids, documents=documents, metadatas=metadatas)
    print(f"  Metadata vectorized and saved: {collection.count()}")


def vectorize_and_save_course_outputs(
    course_name: str,
    chunks_df: pd.DataFrame,
    summaries: dict | None,
    concepts: dict | None,
    metadata_objects: dict,
    chroma_dir: Path,
):
    print("  Starting vectorization and ChromaDB saving...")

    save_chunks_to_chroma(course_name, chunks_df, chroma_dir)

    if summaries:
        summaries_df = build_summaries_df(course_name, summaries)
        debug_semantic(
            f"Vectorizing summaries course={course_name!r}: rows={len(summaries_df)} "
            f"nonempty={int(summaries_df['summary_text'].astype(str).str.strip().astype(bool).sum()) if not summaries_df.empty else 0}"
        )
        save_summaries_to_chroma(course_name, summaries_df, chroma_dir)
    else:
        debug_semantic(f"Vectorization skipped summaries course={course_name!r}: summaries falsy")

    if concepts:
        concepts_df = flatten_concepts_for_chroma(course_name, concepts)
        save_concepts_to_chroma(course_name, concepts_df, chroma_dir)
    else:
        debug_semantic(f"Vectorization skipped concepts course={course_name!r}: concepts falsy")

    metadata_df = build_metadata_documents(course_name, metadata_objects)
    save_metadata_to_chroma(course_name, metadata_df, chroma_dir)

    print("  Vectorization and ChromaDB saving completed.")


# =========================================================
# CHAPTER SUBTOPIC GROUPING + CLUSTERING
# =========================================================

def normalize_subtopic(text: str) -> str:
    text = str(text).strip().lower()
    text = re.sub(r"[_\-]+", " ", text)
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[^\w\s]", "", text)
    return text.strip()


def build_chapter_subtopics_grouped(course_name: str, concepts: dict) -> dict:
    grouped    = {"course_name": course_name, "chapters": []}
    chapter_map: dict = {}

    for relative_path, content in concepts.items():
        chapter = str(content.get("chapter", "")).strip() or "NO_CHAPTER"

        if chapter not in chapter_map:
            chapter_map[chapter] = {"chapter": chapter, "topics": []}

        for topic in content.get("topics", []):
            topic_name = str(topic.get("topic_name", "")).strip()
            subtopics  = [str(s).strip() for s in topic.get("subtopics", []) if str(s).strip()]
            keywords   = [str(k).strip() for k in topic.get("keywords",   []) if str(k).strip()]

            if not topic_name and not subtopics:
                continue

            chapter_map[chapter]["topics"].append({
                "topic_name":    topic_name,
                "subtopics":     subtopics,
                "keywords":      keywords,
                "relative_path": relative_path,
            })

    grouped["chapters"] = sorted(chapter_map.values(), key=lambda x: x["chapter"])
    return grouped


def build_subtopic_clusters(course_name: str, concepts: dict) -> dict:
    clusters_output = {"course_name": course_name, "chapters": []}
    chapter_buckets: dict = {}

    for relative_path, content in concepts.items():
        chapter = str(content.get("chapter", "")).strip() or "NO_CHAPTER"

        if chapter not in chapter_buckets:
            chapter_buckets[chapter] = []

        for topic in content.get("topics", []):
            topic_name = str(topic.get("topic_name", "")).strip()
            keywords   = [str(k).strip() for k in topic.get("keywords", []) if str(k).strip()]

            for subtopic in topic.get("subtopics", []):
                subtopic = str(subtopic).strip()
                if not subtopic:
                    continue
                chapter_buckets[chapter].append({
                    "topic_name":          topic_name,
                    "subtopic_name":       subtopic,
                    "normalized_subtopic": normalize_subtopic(subtopic),
                    "keywords":            keywords,
                    "relative_path":       relative_path,
                })

    for chapter, items in sorted(chapter_buckets.items(), key=lambda x: x[0]):
        used: set         = set()
        chapter_clusters: list = []

        for i, item in enumerate(items):
            if i in used:
                continue

            current_cluster = [item]
            used.add(i)
            item_words = set(item["normalized_subtopic"].split())

            for j, other in enumerate(items):
                if j in used:
                    continue
                other_words = set(other["normalized_subtopic"].split())
                if not item_words or not other_words:
                    continue
                overlap = len(item_words & other_words)
                min_len = min(len(item_words), len(other_words))
                if (
                    item["normalized_subtopic"] == other["normalized_subtopic"]
                    or (min_len > 0 and overlap / min_len >= 0.6)
                ):
                    current_cluster.append(other)
                    used.add(j)

            representative = current_cluster[0]["subtopic_name"]
            chapter_clusters.append({
                "cluster_id":              f"{safe_slug(course_name)}_{safe_slug(chapter)}_{len(chapter_clusters)+1}",
                "chapter":                 chapter,
                "representative_subtopic": representative,
                "members": [
                    {
                        "topic_name":    x["topic_name"],
                        "subtopic_name": x["subtopic_name"],
                        "relative_path": x["relative_path"],
                    }
                    for x in current_cluster
                ],
                "cluster_size": len(current_cluster),
            })

        clusters_output["chapters"].append({"chapter": chapter, "clusters": chapter_clusters})

    return clusters_output


# =========================================================
# MAIN COURSE PIPELINE
# =========================================================

def load_json_file(path: Path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def write_json_debug(path: Path, data, label: str) -> None:
    existed = path.exists()
    old_size = path.stat().st_size if existed else 0
    if isinstance(data, dict):
        shape = f"dict_keys={len(data)}"
    elif isinstance(data, list):
        shape = f"list_len={len(data)}"
    else:
        shape = f"type={type(data).__name__}"
    debug_semantic(
        f"WRITE {label}: path={path} existed={existed} old_size={old_size} {shape}"
    )
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    new_size = path.stat().st_size if path.exists() else 0
    debug_semantic(f"WRITE complete {label}: path={path} new_size={new_size}")


def is_course_already_processed(course_dir: Path) -> bool:
    outputs_dir  = course_dir / "outputs"
    metadata_dir = course_dir / "metadata"
    chroma_dir   = outputs_dir / "chroma_db"

    required_files = [
        outputs_dir / "raw_docs.json",
        outputs_dir / "chapters.csv",
        outputs_dir / "chunks.csv",
        metadata_dir / "course_info.json",
        metadata_dir / "prerequisites.json",
        metadata_dir / "materials_index.json",
        metadata_dir / "topic_schema.json",
    ]

    if RUN_SUMMARIZATION:
        required_files.append(outputs_dir / "chapter_summaries.json")

    if RUN_CONCEPT_EXTRACTION:
        required_files.append(outputs_dir / "chapter_concepts.json")
        required_files.append(outputs_dir / "chapter_subtopics_grouped.json")
        required_files.append(outputs_dir / "subtopic_clusters.json")

    files_exist = all(path.exists() for path in required_files)

    if SAVE_TO_CHROMA:
        return files_exist and chroma_dir.exists()

    return files_exist


def process_course(course_dir: Path):
    course_start = time.time()
    course_name = course_dir.name
    outputs_dir = course_dir / "outputs"
    chroma_dir  = outputs_dir / "chroma_db"

    def print_course_time() -> None:
        elapsed_minutes = (time.time() - course_start) / 60
        print(f"[COURSE TIME] {course_name} took {elapsed_minutes:.2f} min")

    outputs_dir.mkdir(parents=True, exist_ok=True)

    chapter_grouped_file   = outputs_dir / "chapter_subtopics_grouped.json"
    subtopic_clusters_file = outputs_dir / "subtopic_clusters.json"
    concepts_file          = outputs_dir / "chapter_concepts.json"

    # Partial processing: only create missing clustering files
    if not FORCE_REGENERATE:
        if concepts_file.exists() and (
            not chapter_grouped_file.exists() or not subtopic_clusters_file.exists()
        ):
            print(f"\nPartial processing (only clustering) for: {course_name}")
            concepts = load_json_file(concepts_file)
            topic_count, subtopic_count = count_concept_items(concepts)
            debug_semantic(
                f"Partial clustering source concepts course={course_name!r} "
                f"topics={topic_count} subtopics={subtopic_count}"
            )
            write_json_debug(
                chapter_grouped_file,
                build_chapter_subtopics_grouped(course_name, concepts),
                "chapter_grouped_partial",
            )
            write_json_debug(
                subtopic_clusters_file,
                build_subtopic_clusters(course_name, concepts),
                "subtopic_clusters_partial",
            )
            print("  Clustering done ✅ (no reprocessing)")
            print_course_time()
            return

    if not FORCE_REGENERATE and is_course_already_processed(course_dir):
        debug_semantic(
            f"SKIP course={course_name!r}: FORCE_REGENERATE={FORCE_REGENERATE}, "
            "is_course_already_processed=True"
        )
        print(f"\nSkipping already processed course: {course_name} (outputs already exist)")
        print_course_time()
        return

    print(f"\nProcessing course: {course_name}")

    files = collect_supported_sources(course_dir)

    if not files:
        print("  No supported files found. Skipping.")
        print_course_time()
        return

    print(f"  Sources found: {len(files)}")

    # OPT: hash-based change detection — skip unchanged files on re-runs
    changed_files, new_hash_cache = _files_changed_since_last_run(course_dir, files)
    debug_semantic(
        f"Change detection course={course_name!r}: total_files={len(files)} "
        f"changed_files={len(changed_files)} FORCE_REGENERATE={FORCE_REGENERATE}"
    )
    if not changed_files:
        debug_semantic(
            f"SKIP course={course_name!r}: no changed files; semantic outputs left unchanged"
        )
        print(f"  All {len(files)} files unchanged — skipping.")
        print_course_time()
        return
    if len(changed_files) < len(files):
        print(f"  {len(changed_files)}/{len(files)} files changed — processing changed files only.")
        files = changed_files

    docs = load_all_docs(files, course_dir)

    if not docs:
        print("  No readable documents extracted. Skipping.")
        print_course_time()
        return

    run_failed_vision_recovery_pass(course_dir, docs)
    print_vision_counters()

    docs, removed_lines = remove_repeated_page_lines(docs)

    write_json_debug(outputs_dir / "raw_docs.json", docs, "raw_docs")

    write_json_debug(
        outputs_dir / "removed_repeated_lines.json",
        sorted(list(removed_lines)),
        "removed_repeated_lines",
    )

    docs_df = build_docs_df(docs)
    docs_df.to_csv(outputs_dir / "chapters.csv", index=False)
    debug_semantic(
        f"WRITE chapters.csv course={course_name!r}: rows={len(docs_df)} "
        f"nonempty_full_text={int(docs_df['full_text'].astype(str).str.strip().astype(bool).sum())}"
    )

    metadata_objects = write_metadata_files(course_dir, docs_df)

    chunks_df = build_chunks_df(course_name, docs_df)
    chunks_df.to_csv(outputs_dir / "chunks.csv", index=False)
    debug_semantic(f"WRITE chunks.csv course={course_name!r}: rows={len(chunks_df)}")

    print(f"  Chunks created: {len(chunks_df)}")

    summaries      = None
    concepts       = None
    summaries_file = outputs_dir / "chapter_summaries.json"

    if RUN_SUMMARIZATION:
        if summaries_file.exists() and not FORCE_REGENERATE:
            summaries = load_json_file(summaries_file)
            nonempty = sum(bool(str(v).strip()) for v in summaries.values()) if isinstance(summaries, dict) else 0
            debug_semantic(
                f"SKIP summarization course={course_name!r}: existing={summaries_file} "
                f"entries={len(summaries) if isinstance(summaries, dict) else 'n/a'} nonempty={nonempty}"
            )
            print("  Summaries already exist, skipping LLM summarization")
        else:
            summaries = run_summarization(course_name, docs_df)  # OPT: now parallel
            nonempty = sum(bool(str(v).strip()) for v in summaries.values())
            debug_semantic(
                f"Preparing to save summaries course={course_name!r}: "
                f"entries={len(summaries)} nonempty={nonempty} empty={len(summaries) - nonempty}"
            )
            write_json_debug(summaries_file, summaries, "chapter_summaries")
            print("  Summaries saved")

    if RUN_CONCEPT_EXTRACTION:
        if concepts_file.exists() and not FORCE_REGENERATE:
            concepts = load_json_file(concepts_file)
            topic_count, subtopic_count = count_concept_items(concepts)
            debug_semantic(
                f"SKIP concept extraction course={course_name!r}: existing={concepts_file} "
                f"files={len(concepts) if isinstance(concepts, dict) else 'n/a'} "
                f"topics={topic_count} subtopics={subtopic_count}"
            )
            print("  Concepts already exist, skipping LLM concept extraction")
        else:
            concepts = run_concept_extraction(course_name, docs_df)  # OPT: now parallel
            topic_count, subtopic_count = count_concept_items(concepts)
            debug_semantic(
                f"Preparing to save concepts course={course_name!r}: "
                f"files={len(concepts)} topics={topic_count} subtopics={subtopic_count}"
            )
            write_json_debug(concepts_file, concepts, "chapter_concepts")
            print("  Concepts saved")

        chapter_grouped = build_chapter_subtopics_grouped(course_name, concepts)
        debug_semantic(
            f"Preparing to save chapter-grouped subtopics course={course_name!r}: "
            f"chapters={len(chapter_grouped.get('chapters', []))}"
        )
        write_json_debug(chapter_grouped_file, chapter_grouped, "chapter_grouped")
        print("  Chapter-grouped subtopics saved")

        subtopic_clusters = build_subtopic_clusters(course_name, concepts)
        debug_semantic(
            f"Preparing to save subtopic clusters course={course_name!r}: "
            f"chapters={len(subtopic_clusters.get('chapters', []))}"
        )
        write_json_debug(subtopic_clusters_file, subtopic_clusters, "subtopic_clusters")
        print("  Subtopic clusters saved")

    if SAVE_TO_CHROMA:
        vectorize_and_save_course_outputs(
            course_name=course_name,
            chunks_df=chunks_df,
            summaries=summaries,
            concepts=concepts,
            metadata_objects=metadata_objects,
            chroma_dir=chroma_dir,
        )

    # OPT: save hash cache only after successful completion
    _save_hash_cache(course_dir, new_hash_cache)

    print("  Metadata written:", list(metadata_objects.keys()))
    print("  Done")
    print_course_time()


def process_all_courses(courses_root: Path):
    course_folders = collect_course_folders(courses_root)

    if not course_folders:
        print("No course folders found.")
        return

    print(f"Found {len(course_folders)} course folders.\n")

    if TEST_COURSE_NAME:
        test_course = courses_root / TEST_COURSE_NAME
        if not test_course.exists():
            raise FileNotFoundError(f"Test course folder not found: {test_course}")
        process_course(test_course)
        print("\nTest course processed successfully.")
        print_preprocessing_statistics()
        return

    if MAX_COURSE_WORKERS > 1:
        with ThreadPoolExecutor(max_workers=MAX_COURSE_WORKERS) as pool:
            futures = {pool.submit(process_course, course_dir): course_dir.name for course_dir in course_folders}
            for future in as_completed(futures):
                name = futures[future]
                try:
                    future.result()
                except Exception as e:
                    print(f"\n[ERROR] Course '{name}' failed: {e}")
    else:
        for course_dir in course_folders:
            try:
                process_course(course_dir)
            except Exception as e:
                print(f"\n[ERROR] Course '{course_dir.name}' failed: {e}")

    print("\nAll courses processed successfully.")
    print_preprocessing_statistics()


# =========================================================
# METADATA-ONLY REFRESH
# =========================================================

def refresh_metadata_for_all_courses(courses_root: Path = COURSES_DIR) -> None:
    """
    Re-write the four lightweight metadata JSON files for every already-processed
    course WITHOUT re-running any LLM calls or ChromaDB operations.

    Use this after changing COURSE_METADATA_MAP (e.g. fixing prerequisites from
    course codes to course names) so that existing courses pick up the changes
    without requiring a full FORCE_REGENERATE run.

    Only touches:
      metadata/course_info.json
      metadata/prerequisites.json
      metadata/materials_index.json
      metadata/topic_schema.json

    Skips courses that have never been processed (no outputs/ folder).
    """
    course_folders = collect_course_folders(courses_root)
    updated = 0
    skipped = 0

    for course_dir in course_folders:
        outputs_dir  = course_dir / "outputs"
        metadata_dir = course_dir / "metadata"

        if not outputs_dir.exists() or not metadata_dir.exists():
            print(f"  SKIP (not yet processed): {course_dir.name}")
            skipped += 1
            continue

        chapters_csv = outputs_dir / "chapters.csv"
        if not chapters_csv.exists():
            print(f"  SKIP (no chapters.csv): {course_dir.name}")
            skipped += 1
            continue

        try:
            docs_df = pd.read_csv(chapters_csv)
            write_metadata_files(course_dir, docs_df)
            print(f"  REFRESHED metadata: {course_dir.name}")
            updated += 1
        except Exception as e:
            print(f"  ERROR refreshing {course_dir.name}: {e}")

    print(f"\nMetadata refresh complete — {updated} updated, {skipped} skipped.")


# =========================================================
# RUN
# =========================================================

if __name__ == "__main__":
    # To update only metadata (prerequisites, course info) for already-processed
    # courses after changing COURSE_METADATA_MAP, run:
    #   refresh_metadata_for_all_courses()
    #
    # To process new courses only (skip already-processed ones), keep
    # FORCE_REGENERATE = False and run:
    #   process_all_courses(COURSES_DIR)
    #
    # To reprocess everything from scratch, set FORCE_REGENERATE = True first.

    print_startup_rate_limit_config()
    process_all_courses(COURSES_DIR)
